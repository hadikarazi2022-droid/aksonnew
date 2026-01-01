# base.py (Farooqi AI Widget — Exact UI match + X / Minimise, no other changes!)
from PyQt6.QtWidgets import QLineEdit
from PyQt6 import QtWidgets, QtCore
from PyQt6.QtWidgets import QApplication, QMainWindow, QTextEdit, QPushButton, QHBoxLayout
from PyQt6.QtGui import QKeySequence, QAction, QIcon
from PyQt6.QtCore import QStandardPaths
import sys
import platform
import mss
import cv2
from PyQt6.QtWidgets import QTextBrowser, QScrollBar

from PyQt6.QtGui import QTextOption
import csv
from PyQt6 import QtGui
import numpy as np
import openai
from pynput import keyboard
import os, json
# Temporarily disabled for standalone operation
# from login import LoginDialog
# from akson.config import OPENAI_API_KEY, validate_config
# from request_limit import get_user_data, can_make_request, increment_request_count
import requests
# from firebase_admin import auth
from PyQt6.QtWidgets import QApplication
from PyQt6.QtCore import pyqtSlot, QTimer
from PyQt6.QtWidgets import QDialog, QLabel, QVBoxLayout, QDialogButtonBox, QPushButton, QToolButton, QCheckBox
from PyQt6.QtWebEngineWidgets import QWebEngineView
import threading
import time
import webbrowser
import subprocess
import pytesseract
from PIL import Image
from collections import deque
from pynput.keyboard import Key
import genanki
from PyQt6.QtPrintSupport import QPrinter
from PyQt6 import QtGui as _QtGui  # already imported QtGui above, this is just to reference QTextDocument
from typing import Optional

# Slides import deps
try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import fitz  # PyMuPDF for PDFs
except Exception:
    fitz = None



# --- SUMMARY POST-PROCESSOR: compress, split, symbols, cap length ---


def markdown_to_html(text):
    import html, re
    text = html.escape(text)

    # optional small headers if they appear
    text = re.sub(r'^##\s+(.*)$',
                  r"<span style='font-size:15px;font-weight:600;display:block;margin:6px 0 4px;'>\1</span>",
                  text, flags=re.MULTILINE)
    text = re.sub(r'^###\s+(.*)$',
                  r"<span style='font-size:14px;font-weight:600;display:block;margin:4px 0 2px;'>\1</span>",
                  text, flags=re.MULTILINE)

    # bold/italic/underline
    text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)
    text = re.sub(r'\*(.+?)\*', r'<i>\1</i>', text)
    text = re.sub(r'~~(.+?)~~', r'<u>\1</u>', text)

    # bullets → use larger font + relaxed line-height
    text = re.sub(r'(?m)^\-\s+(.*)$',
                  r"<div style='margin:8.5px 0;line-height:1.1;font-size:14px;'>• \1</div>",
                  text)

    # drop remaining newlines
    text = re.sub(r'\n+', '', text).strip()

    # wrap everything in a global style block with width constraints
    return f"<div style='font-size:14.5px; line-height:1.1; font-family:-apple-system,BlinkMacSystemFont,Segoe UI,Roboto,sans-serif; max-width:100%; word-wrap:break-word; overflow-wrap:break-word;'>{text}</div>"




# Initialize OpenAI API Key from config
try:
    from akson.config import OPENAI_API_KEY
    openai.api_key = OPENAI_API_KEY
except ImportError:
    # Fallback for when running without akson package
    import os
    openai.api_key = os.getenv("OPENAI_API_KEY")

def load_or_prompt_user():
    # Standalone mode - bypass login
    if os.path.exists("user.json"):
        with open("user.json", "r") as f:
            user_data = json.load(f)
            if "custom_prompt" not in user_data:
                user_data["custom_prompt"] = ""
            return user_data
    
    # Return default user for standalone operation
    return {
        "email": "standalone@localhost",
        "custom_prompt": "",
        "consent_version": "2025-08-31",
        "is_premium": True  # No limits in standalone mode
    }


 # <-- add this import at top


POLICY_VERSION = "2025-08-31"

POLICY_HTML = """
<h2 style="margin:0 0 8px 0;">Akson — Privacy Notice & Terms</h2>
<p style="margin:0 0 12px 0; font-size:13px;">
By using Akson you agree to the Privacy Policy and Terms below. You must scroll to the end to enable "I agree".
</p>

<h3 style="margin:14px 0 6px 0;">How your data is processed</h3>
<ul style="margin:6px 0 12px 18px; line-height:1.3;">
  <li>Akson helps you study by taking text from your screen (via screenshots or copied text) and processing it with AI to give instant answers.</li>
  <li><b>We do not store, see, or log</b> screen captures or selected text. They are processed on your device and the extracted text is sent directly to OpenAI for AI processing.</li>
  <li>No humans view your content.</li>
  <li>We keep only your email and subscription status for login/billing.</li>
  <li>You can request account deletion any time by emailing us.</li>
  <li>Do not capture sensitive/confidential data (e.g., patient records, passwords).</li>
  <li>By using Akson, you also agree to OpenAI's API Terms of Use.</li>
</ul>

<hr style="border:none;border-top:1px solid #333; margin:12px 0;">

<h3 style="margin:14px 0 6px 0;">Privacy Policy</h3>
<p style="margin:0 0 6px 0;">Effective Date: 31 August 2025<br>Last Updated: 31 August 2025<br>
Data Protection Officer (DPO): Taha Fareed Farooqi — <a href="mailto:tahaffarooqi@gmail.com">tahaffarooqi@gmail.com</a></p>

<h4>1. Who We Are</h4>
<p>Akson is a study assistant that uses OCR and AI to generate educational content.</p>

<h4>2. What Data We Collect</h4>
<p><b>2.1 Account Information</b>: Email; subscription status.<br>
<b>2.2 Usage Information</b>: Daily request counts; error logs (technical only).<br>
<b>2.3 OCR, Clipboard & AI Data</b>: Screenshots are temporary on-device; extracted text goes to OpenAI; we do not store it.<br>
<b>2.4 Keyboard Shortcuts & Clipboard Access</b>: Listens only for Ctrl+C and Right/Down arrows. No other keystrokes are read or stored.</p>

<h4>3. How We Use Data</h4>
<p>Login/billing; generate AI outputs; enforce limits; reliability and fixes.</p>

<h4>4. Legal Basis (GDPR)</h4>
<p>Consent (temporary processing of captured text) and legitimate interests (minimal account metrics).</p>

<h4>5. Data Sharing</h4>
<p>OpenAI (processing), Firebase (auth), payment processor (billing). No advertisers; no data selling.</p>

<h4>6. Data Retention</h4>
<p>Account data persists while active; usage counts reset daily; OCR/AI content not stored.</p>

<h4>7. Your Rights</h4>
<p>Access, correction, deletion, portability, and complaints to your DPA.</p>

<h4>8. Security</h4>
<p>Firebase Auth; keys protected/excluded from builds; no storage of OCR content.</p>

<h4>9. Children's Privacy</h4>
<p>Not for under 18s.</p>

<h4>10. Educational Use Only (Disclaimer & Liability)</h4>
<p>Educational tool; do not use for clinical decisions; verify with official sources. You are responsible for captured content. We aren't liable for misuse; information may have errors; not medical/legal advice.</p>

<h4>11. Updates</h4>
<p>We may update this policy; we'll notify when required by law.</p>

<h4>12. Contact</h4>
<p>DPO: Taha Fareed Farooqi — <a href="mailto:tahaffarooqi@gmail.com">tahaffarooqi@gmail.com</a></p>

<h3 style="margin:14px 0 6px 0;">Terms of Use (Akson)</h3>
<p>Effective Date: 31 July 2025<br>Last Updated: 31 July 2025</p>

<ol style="margin:6px 0 12px 18px; line-height:1.3;">
  <li><b>Acceptance</b>: Using the app binds you to these Terms and the Privacy Policy.</li>
  <li><b>Educational Purpose</b>: Not a medical device; not for clinical decisions.</li>
  <li><b>User Responsibilities</b>: Lawful use; right to capture; avoid sensitive/confidential data; keep account secure.</li>
  <li><b>Data Handling & Third Parties</b>: OCR → text to OpenAI; no storage of content; Firebase for auth; payments via third party; you also agree to OpenAI's terms.</li>
  <li><b>Subscription</b>: Free daily request limits; Premium lifts limits; refunds per provider/law.</li>
  <li><b>IP</b>: No copying/reverse engineering of Akson.</li>
  <li><b>Warranties</b>: "As is"; no guarantee of accuracy or uptime.</li>
  <li><b>Liability</b>: Max liability limited to what you paid in last 12 months.</li>
  <li><b>Termination</b>: We may suspend/terminate for violations/abuse.</li>
  <li><b>Governing Law</b>: Republic of Ireland.</li>
  <li><b>Changes</b>: Terms may change; continued use = acceptance.</li>
  <li><b>Contact</b>: <a href="mailto:tahaffarooqi@gmail.com">tahaffarooqi@gmail.com</a></li>
</ol>

<p><b>Key Legal Protections for You:</b> We don't store/view your screen/clipboard; you decide what to capture; GDPR compliant; AI processing by OpenAI.</p>

<p style="margin:12px 0 0 0;">
OpenAI API Terms of Use: <a href="https://openai.com/policies/terms-of-use">https://openai.com/policies/terms-of-use</a>
</p>
"""



def _save_user_json_merge(new_fields: dict):
    data = {}
    if os.path.exists("user.json"):
        try:
            with open("user.json", "r") as f:
                data = json.load(f) or {}
        except Exception:
            data = {}
    data.update(new_fields or {})
    try:
        with open("user.json", "w") as f:
            json.dump(data, f)
    except Exception as e:
        print(f"[Consent] Failed saving user.json: {e}")

class ConsentDialog(QDialog):
    """
    Mandatory consent dialog shown post-login and before app launch.
    Requires user to scroll to bottom before enabling the 'I agree' checkbox.
    """
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Privacy & Terms — Consent Required")
        self.setModal(True)
        self.setFixedSize(520, 560)
        self.setStyleSheet("""
            QDialog { background:#1e1e1e; color:#dcdcdc; border-radius:8px; }
            QLabel  { color:#dcdcdc; font-size:13.5px; }
            QTextBrowser {
                background:#141416; color:#e6e6e6;
                border:0.5px solid #2a2a2e; border-radius:8px;
                padding:8px; font-size:13px;
            }
            QCheckBox { color:#dcdcdc; font-size:13.5px; }
            QPushButton {
                background: transparent; color: #dcdcdc;
                padding: 6px 12px; border-radius: 8px; font-size: 13.5px; font-weight:600;
                border: 0.5px solid #2a2a2e;
            }
            QPushButton:hover { background-color:#1e1e1e; border:0.5px solid #7f7fff; color:#d0caff; }
            QDialogButtonBox { padding-top:6px; }
        """)

        v = QVBoxLayout(self)
        v.setContentsMargins(12,12,12,12)
        v.setSpacing(8)

        info = QLabel("Please review the policy. Scroll to the bottom to enable agreement.")
        info.setWordWrap(True)
        v.addWidget(info)

        self.viewer = QTextBrowser()
        self.viewer.setOpenExternalLinks(True)
        self.viewer.setHtml(POLICY_HTML)
        v.addWidget(self.viewer, 1)

        self.chk = QCheckBox("I have read and agree to the Privacy Policy and Terms.")
        self.chk.setEnabled(False)  # enabled only after full scroll
        v.addWidget(self.chk)

        self.btns = QDialogButtonBox()
        self.btn_agree = QPushButton("Agree & Continue")
        self.btn_cancel = QPushButton("Cancel")
        self.btn_agree.setEnabled(False)
        self.btns.addButton(self.btn_cancel, QDialogButtonBox.ButtonRole.RejectRole)
        self.btns.addButton(self.btn_agree, QDialogButtonBox.ButtonRole.AcceptRole)
        v.addWidget(self.btns)

        # wiring
        self.viewer.verticalScrollBar().valueChanged.connect(self._on_scroll)
        self.chk.toggled.connect(lambda s: self.btn_agree.setEnabled(bool(s)))
        self.btn_cancel.clicked.connect(self.reject)
        self.btn_agree.clicked.connect(self._on_accept)

    def _on_scroll(self, _):
        sb: QScrollBar = self.viewer.verticalScrollBar()
        at_bottom = sb.value() >= (sb.maximum() - 4)
        if at_bottom:
            self.chk.setEnabled(True)

    def _on_accept(self):
        if not self.chk.isChecked():
            return
        self.accept()


def sync_premium_status(user_data):
    # Standalone mode - always premium, no limits
    user_data["is_premium"] = True
    print("Standalone mode: premium features enabled")




# Sync premium status from users.json (written by server)



class SettingsDialog(QDialog):
    def __init__(self, user_data, parent=None):
        super().__init__(parent)
        self.user_data = user_data
        self.keyboard_enabled = True  # default to True
        self.setWindowOpacity(0.96)

        self.setWindowTitle("Settings")
        self.setFixedSize(315, 500)
        self.setStyleSheet("""
    QDialog {
        background-color: #1e1e1e;
        color: #dcdcdc;
        font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
        font-size: 13.5px;
        border-radius: 8px;
    }
    QLabel {
        padding: 6px 4px;
        color: #dcdcdc;
        font-size: 13.5px;
    }
    QPushButton {
        background: transparent;
        color: #dcdcdc;
        padding: 6px 12px;
        border-radius: 8px;
        font-size: 13.5px;
        font-weight: 600;
        border: 0.5px solid #2a2a2e;
    }
    QPushButton:hover {
        background-color: #1e1e1e;
        border: 0.5px solid #7f7fff;
        color: #d0caff;
    }
""")

        
                

        

        layout = QVBoxLayout()

        today = time.strftime('%Y-%m-%d')
        request_count = self.user_data.get('requests', {}).get(today, 0)

        email_label = QLabel(f"Email: {self.user_data.get('email', 'N/A')}")
        plan_label = QLabel(f"Plan: {'Premium' if self.user_data.get('is_premium') else 'Free'}")
        count_label = QLabel(f"Today's Requests: {request_count}")
        for label in [email_label, plan_label, count_label]:
            label.setStyleSheet("color: #cccccc; font-size: 13.5px; padding: 4px;")


        layout.addWidget(email_label)
        layout.addWidget(plan_label)
        layout.addWidget(count_label)

        upgrade_button = QPushButton("Upgrade Plan")
        upgrade_button.clicked.connect(self.open_payment_page)
        layout.addWidget(upgrade_button)

        logout_button = QPushButton("Log Out")
        logout_button.clicked.connect(self.logout)
        layout.addWidget(logout_button)

        refresh_button = QPushButton("Refresh Status")
        refresh_button.clicked.connect(self.refresh_status)
        layout.addWidget(refresh_button)

        self.down_arrow_checkbox = QCheckBox("Enable Down Arrow Shortcut")
        self.down_arrow_checkbox.setStyleSheet("color: #dcdcdc;")
        self.down_arrow_checkbox.setChecked(self.parent().keyboard_enabled)
        layout.addWidget(self.down_arrow_checkbox)

        # Label for clarity
        # Label for clarity
        prompt_label = QLabel("Your Custom AI Instructions:")
        prompt_label.setStyleSheet("color: #dcdcdc; font-size: 13.5px; padding: 4px;")
        layout.addWidget(prompt_label)

        # Prompt input box (pre-filled)
        self.prompt_input = QtWidgets.QTextEdit()
        self.prompt_input.setPlaceholderText("Add custom instructions here (optional)...")
        self.prompt_input.setFixedHeight(80)
        self.prompt_input.setStyleSheet("""
            QTextEdit {
                background-color: #1e1e1e;
                color: #dcdcdc;
                border: 0.5px solid #2a2a2e;
                border-radius: 8px;
                font-size: 13.5px;
                padding: 6px;
            }
        """)
        layout.addWidget(self.prompt_input)

        # === Load existing prompt from user_data ===
        existing_prompt = self.user_data.get("custom_prompt", "")
        print(f"[DEBUG] custom_prompt in user_data: {self.user_data.get('custom_prompt')}")

        if existing_prompt:
            self.prompt_input.setPlainText(existing_prompt)


        reset_button = QPushButton("Reset Custom Prompt")
        reset_button.clicked.connect(lambda: self.prompt_input.clear())

        layout.addWidget(reset_button)

        buttons = QDialogButtonBox(QDialogButtonBox.StandardButton.Ok)
        buttons.accepted.connect(self.accept)
        layout.addWidget(buttons)

        self.setLayout(layout)

        if self.user_data.get("is_premium"):
            upgrade_button.setEnabled(False)
            upgrade_button.setText("Already Premium ✅")

    def logout(self):
        confirm = QtWidgets.QMessageBox.question(
            self,
            "Confirm Logout",
            "Are you sure you want to log out?"
        )
        if confirm == QtWidgets.QMessageBox.StandardButton.Yes:
            # remove saved credentials
            try:
                if os.path.exists("user.json"):
                    os.remove("user.json")
            except Exception as e:
                QtWidgets.QMessageBox.critical(self, "Error", f"Failed to log out: {e}")
                return

            QtWidgets.QMessageBox.information(self, "Logged Out", "You have been logged out.")
            # exit the Qt event loop with code 100 to signal a logout
            QtWidgets.QApplication.exit(100)
  # Return special code to trigger restart
  # Trigger full restart



    def open_payment_page(self):
        email = self.user_data.get("email")
        if not email:
            QtWidgets.QMessageBox.critical(self, "Error", "No email found in your user profile.")
            return
        url = f"https://farooqi-payments.onrender.com/create-checkout-session?email={email}"
        webbrowser.open(url)

    def refresh_status(self):
        sync_premium_status(self.user_data)
        QtWidgets.QMessageBox.information(self, "Status Updated", "Premium status refreshed. Close and reopen Settings to see changes.")




class CompactBar(QtWidgets.QWidget):
    # UI-thread signals
    result_ready = QtCore.pyqtSignal(str)
    error_ready  = QtCore.pyqtSignal(str)

    def __init__(self, main_window: "MyWindow"):
        super().__init__(None)
        self.main_window = main_window
        self._answer_displayed = False  # tracks if the box currently shows an answer


        # Always-on-top window (stays above other apps)
        self.setWindowFlags(
            QtCore.Qt.WindowType.FramelessWindowHint |
            QtCore.Qt.WindowType.WindowStaysOnTopHint |
            QtCore.Qt.WindowType.Window
        )
        self.setAttribute(QtCore.Qt.WidgetAttribute.WA_TranslucentBackground, True)
        self.setFocusPolicy(QtCore.Qt.FocusPolicy.StrongFocus)
        self.setWindowOpacity(0.85)

        # ---------- shell (black container) ----------
        shell = QtWidgets.QFrame(self)
        shell.setObjectName("miniShell")
        shell.setStyleSheet("""
            #miniShell { background-color: #1e1e1e; border: 0.5px solid #2a2a2e; border-radius: 10px; }
            QLabel { color: #cccccc; font-size: 12.5px; padding: 0px; }
            QTextEdit {
                background-color: #1e1e1e; color: #dcdcdc; border: 0.5px solid #2a2a2e;
                border-radius: 8px; padding: 4px 8px; font-size: 12.5px; min-width: 220px;
            }
            QPushButton {
                background: transparent; color: #dcdcdc; padding: 4px 10px; border-radius: 8px;
                font-size: 12.5px; font-weight: 600; border: 0.5px solid #2a2a2e;
            }
            QPushButton:hover { background-color: #1e1e1e; border: 0.5px solid #7f7fff; color: #d0caff; }
        """)
        row = QtWidgets.QHBoxLayout(shell)
        row.setContentsMargins(12, 4, 12, 4)   # tight so the black matches the controls
        row.setSpacing(4)

        # mode chip (same look)



        # ---------- input (keeps old look; now auto-height) ----------
        self.ask_input = QtWidgets.QTextEdit()
        self.ask_input.setAcceptRichText(False)
        self.ask_input.setWordWrapMode(QtGui.QTextOption.WrapMode.WrapAtWordBoundaryOrAnywhere)
        self.ask_input.setPlaceholderText("Type & Enter…")
        self.ask_input.setVerticalScrollBarPolicy(QtCore.Qt.ScrollBarPolicy.ScrollBarAlwaysOff)
        self.ask_input.setSizePolicy(QtWidgets.QSizePolicy.Policy.Expanding,
                                     QtWidgets.QSizePolicy.Policy.Fixed)
        self.ask_input.document().setDocumentMargin(2)  # trims inner top/bottom padding
        self.ask_input.setFrameShape(QtWidgets.QFrame.Shape.NoFrame)

        self._min_input_h = 28
        self._max_input_h = 120
        self.ask_input.setFixedHeight(self._min_input_h)

        # Submit on Enter; Shift+Enter = newline
        self.ask_input.installEventFilter(self)
        # Auto-grow when answer wraps
        self.ask_input.textChanged.connect(self._adjust_height_to_content)

        # ---------- buttons (Card → X; same size as input) ----------
        self.go_btn = QtWidgets.QPushButton("Go")
        self.go_btn.setToolTip("Summarize screen")
        self.go_btn.clicked.connect(self.main_window.manual_summarize)

        self.close_btn = QtWidgets.QPushButton("✕")   # replaces Card
        self.close_btn.setFixedWidth(28)
        self.close_btn.setToolTip("Close mini")
        self.close_btn.clicked.connect(self._expand)  # exit mini (restore full)

        self.expand_btn = QtWidgets.QPushButton("⤢")
        self.expand_btn.setFixedWidth(28)
        self.expand_btn.setToolTip("Expand")
        self.expand_btn.clicked.connect(self._expand)

        # match heights with the input
        for b in (self.go_btn, self.close_btn, self.expand_btn):
            b.setFixedHeight(self._min_input_h)

        # layout order unchanged (except Card→X)
   
        row.addSpacing(6)
        row.addWidget(self.ask_input)
        row.addWidget(self.go_btn)
        row.addWidget(self.close_btn)
        row.addWidget(self.expand_btn)

        outer = QtWidgets.QHBoxLayout(self)
        outer.setContentsMargins(0, 0, 0, 0)
        outer.addWidget(shell)

        # signals
        self.result_ready.connect(self._apply_answer)
        self.error_ready.connect(self._apply_error)

    # ---------- keep old public API names ----------
    def update_from_main(self): return self._update_from_main()
    def position_top_center(self): return self._position_top_center()
    def adjust_and_reposition(self): return self._adjust_and_reposition()

    # ---------- window behavior (unchanged feel) ----------
    def showEvent(self, e: QtGui.QShowEvent) -> None:
        self._answer_displayed = False

        super().showEvent(e)
        self._update_from_main()
        self._adjust_and_reposition()
        QtCore.QTimer.singleShot(0, self._focus_input)

    def eventFilter(self, obj, event):
        if obj is self.ask_input:
            if event.type() == QtCore.QEvent.Type.KeyPress:
                # If an answer is currently shown, clear it on the first keystroke
                if self._answer_displayed:
                    self._answer_displayed = False
                    self.ask_input.clear()
                    # fall through — this same keypress will now go into an empty box

                if event.key() in (QtCore.Qt.Key.Key_Return, QtCore.Qt.Key.Key_Enter):
                    if not (event.modifiers() & (QtCore.Qt.KeyboardModifier.ShiftModifier |
                                                QtCore.Qt.KeyboardModifier.ControlModifier |
                                                QtCore.Qt.KeyboardModifier.MetaModifier |
                                                QtCore.Qt.KeyboardModifier.AltModifier)):
                        self._ask_from_mini()
                        return True
        return super().eventFilter(obj, event)

    
    def explain_clipboard(self, text: str):
        """
        Mini-mode: on Ctrl+C, show a short but complete explanation/answer of the selected text.
        """
        if not text:
            return

        # UI: show progress
        self.ask_input.setEnabled(False)
        self.ask_input.setPlainText("Working…")
        self._answer_displayed = False

        self._adjust_height_to_content()

        def worker():
            try:
                # Optional recent context from main window
                try:
                    context = getattr(self.main_window, "latest_context", "") or ""
                except Exception:
                    context = ""

                # Same universal mini-mode prompt
                sys_msg = (
                    "MINI MODE — Short but complete clinical answer.\n"
                    "• Answer ONLY the facet implied/asked; do not add other facets. GIVE AT MOST 20-25 WORDS. KEEP IT AS CONCISE BUT RELVANT AS POSSIBLE. give bullet points where possible\n"
        
                    "• If the request expects a list (e.g., treatment/antibiotics/causes/investigations/risk factors/complications), "
                    "return 3–6 SPECIFIC named items, comma-separated. Avoid generic category words alone.\n"
                    "• Otherwise, give 1–2 short sentences with key specifics only (names/numbers) if essential.\n"
                    "• If the input is not a question (a term/phrase), briefly define/explain it for a medical student.\n"
                    "• Use recent context ONLY to disambiguate under-specified input; otherwise answer directly.\n"
                    "• No preamble, no headings, no mention of context."
                )

                user_msg = (
                    f"Selected text:\n{text}\n\n"
                    f"Recent context (optional):\n{context}"
                )

                resp = openai.ChatCompletion.create(
                    model="gpt-4o-mini",
                    messages=[{"role": "system", "content": sys_msg},
                            {"role": "user",   "content": user_msg}],
                    temperature=0.2,
                )
                ans = resp["choices"][0]["message"]["content"].strip()

                self.result_ready.emit(ans)

            except Exception as e:
                print("[MiniExplainClipboard] error:", e)
                self.error_ready.emit("Error")

        threading.Thread(target=worker, daemon=True).start()


    def mousePressEvent(self, event: QtGui.QMouseEvent):
        # Only focus; do NOT expand unless the ⤢ button is pressed
        if event.button() == QtCore.Qt.MouseButton.LeftButton:
            self._focus_input()
        super().mousePressEvent(event)

    def _focus_input(self):
        self.activateWindow()
        self.raise_()
        self.ask_input.setFocus(QtCore.Qt.FocusReason.ActiveWindowFocusReason)

    def _expand(self):
        self.hide()
        self.main_window.exit_compact_mode()

    def _update_from_main(self):
        try:
            self.mode_label.setText(self.main_window.mode_selector.currentText())
        except Exception:
            pass

    def _position_top_center(self):
        screen = QtGui.QGuiApplication.primaryScreen()
        geo = screen.availableGeometry() if screen else QtCore.QRect(0, 0, 1440, 900)
        x = geo.x() + (geo.width() - self.width()) // 2
        y = geo.y() + 8
        self.move(x, y)

    def _adjust_and_reposition(self):
        self.adjustSize()
        self._position_top_center()

    def _adjust_height_to_content(self):
        doc = self.ask_input.document()
        doc.setTextWidth(self.ask_input.viewport().width())
        h = int(doc.size().height()) + 10  # padding
        h = max(self._min_input_h, min(h, self._max_input_h))
        if self.ask_input.height() != h:
            self.ask_input.setFixedHeight(h)
            self._adjust_and_reposition()


    def _wants_specific_list(self, q: str) -> bool:
        ql = q.lower()
        triggers = (
            "treatment", "tx", "management", "mx",
            "antibiotic", "antibiotics", "drug", "medication",
            "cause", "causes", "etiology", "aetiology",
            "risk factor", "risk factors",
            "diagnosis", "dx", "differential", "differentials",
            "investigation", "investigations", "tests", "labs", "imaging",
            "signs", "features", "findings", "complication", "complications",
            "first line", "1st line", "second line", "2nd line"
        )
        return any(t in ql for t in triggers)

    def _is_generic_reply(self, ans: str, q: str) -> bool:
        a = ans.strip().lower()
        # If it's just a category word/phrase and not a list → generic
        generic_terms = {
            "antibiotic", "antibiotics", "antiviral", "antivirals", "supportive care",
            "treatment", "management", "therapy", "medication", "drugs",
            "diagnosis", "investigations", "tests", "labs", "imaging",
            "causes", "cause", "aetiology", "etiology",
            "complications", "symptoms", "signs", "features", "differentials"
        }
        only_words = set(a.replace(",", " ").replace(";", " ").split())
        # If it didn't produce commas/semicolons and is short and overlaps category terms → generic
        looks_like_list = ("," in a) or (";" in a) or ("•" in a) or (" - " in a)
        short = len(only_words) <= 4
        overlaps = any(g in a for g in generic_terms)
        return (self._wants_specific_list(q) and not looks_like_list and overlaps) or (short and overlaps)

    # ---------- mini ask: screenshot + OCR + brief answer ----------
    def _ask_from_mini(self):
        q = self.ask_input.toPlainText().strip()
        if not q:
            return

        # UI: show progress
        self.ask_input.setEnabled(False)
        self.ask_input.setPlainText("Working…")
        self._answer_displayed = False

        self._adjust_height_to_content()

        def worker():
            try:
                # Screenshot WHOLE desktop (use [1] for primary-only if you prefer)
                with mss.mss() as sct:
                    shot = sct.grab(sct.monitors[0])
                img = np.array(shot)
                img = cv2.cvtColor(img, cv2.COLOR_BGRA2BGR)

                # OCR of the screen
                ocr_text = pytesseract.image_to_string(Image.fromarray(img)).strip()

                # Recent context from main window (optional)
                try:
                    context = getattr(self.main_window, "latest_context", "") or ""
                except Exception:
                    context = ""

                # Universal mini-mode prompt: short but complete, facet-locked, specific when lists are expected
                sys_msg = (
                    "MINI MODE — Short but complete clinical answer.\n"

       
                    "• If the question asks for a list (e.g., causes, treatments, antibiotics, complications, risk factors, investigations): give 3–6 named items in bullet points. No generic categories. Always specify drug names if relevant. "
                    "• Do not add headings, preambles, or filler words. Output only the answer.  "
                    "• Answer ONLY the facet asked by the question; do not add other facets.\n"
    
                    "• Use OCR/recent context ONLY to infer the subject when the question is under-specified; "
                    "if the question names the topic, ignore mismatched context and answer directly.\n"
                    "• No preamble, no headings, no mention of OCR/context."
                    "IF MENTIONING MEDICATIONS ACTUALLY SPECIFY WHICH MEDICATIONS!!!"
                )

                user_msg = (
                    f"OCR/context:\n{ocr_text}\n\n"
                    f"Recent extracted context:\n{context}\n\n"
                    f"Question: {q}"
                )

                resp = openai.ChatCompletion.create(
                    model="gpt-4o-mini",
                    messages=[{"role": "system", "content": sys_msg},
                            {"role": "user",   "content": user_msg}],
                    temperature=0.5,
                )
                ans = resp["choices"][0]["message"]["content"].strip()

                self.result_ready.emit(ans)

            except Exception as e:
                print("[MiniAsk] error:", e)
                self.error_ready.emit("Error")

        threading.Thread(target=worker, daemon=True).start()


    # ---------- UI-thread slots ----------


    @QtCore.pyqtSlot(str)
    @QtCore.pyqtSlot(str)
    def _apply_answer(self, ans: str):
        self.ask_input.setEnabled(True)
        self.ask_input.setPlainText(ans)
        self._answer_displayed = True   # <— add this
        self._adjust_height_to_content()
        self.ask_input.moveCursor(QtGui.QTextCursor.MoveOperation.Start)
        self.ask_input.verticalScrollBar().setValue(0)
        self._focus_input()

    @QtCore.pyqtSlot(str)
    def _apply_error(self, msg: str):
        self.ask_input.setEnabled(True)
        self.ask_input.setPlainText(msg)
        self._answer_displayed = True   # <— add this
        self._adjust_height_to_content()
        self.ask_input.moveCursor(QtGui.QTextCursor.MoveOperation.Start)
        self.ask_input.verticalScrollBar().setValue(0)
        self._focus_input()




# ── Simple chooser dialog ─────────────────────────────────────────
class ImportFilesDialog(QtWidgets.QDialog):
    """
    Small page to pick a file, add an optional instruction, and choose actions.
    Returns: (path, instruction, do_summary, do_flashcards, do_quiz)
    """
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Import Files")
        self.setModal(True)
        self.setFixedSize(360, 380)
        self.setStyleSheet("""
            QDialog { background:#1e1e1e; color:white; border-radius:8px; }
            QLabel { color:#ddd; font-size:12px; }
            QLineEdit, QTextEdit {
                background:#1a1a1d; color:white; border:0.5px solid #2a2a2e; border-radius:8px;
                font-size:13.5px; padding:6px;
            }
            QPushButton {
                background: transparent; color: #dcdcdc; border: 0.5px solid # 2a2a2e; border-radius: 8px;
                padding: 6px 12px; font-weight:200; font-size:10px;
            }
            QPushButton:hover { background-color: #1e1e1e; border: 0.5px solid #7f7fff; color:#d0caff; }
            QCheckBox { color:#ddd; font-size:13.5px; }
        """)

        layout = QtWidgets.QVBoxLayout(self)
        layout.setContentsMargins(12,12,12,12)
        layout.setSpacing(8)

        # File picker
        layout.addWidget(QtWidgets.QLabel("Choose a file (.pptx or .pdf):"))
        file_row = QtWidgets.QHBoxLayout()
        self.path_edit = QtWidgets.QLineEdit()
        self.path_edit.setReadOnly(True)
        browse = QtWidgets.QPushButton("Browse")
        browse.clicked.connect(self._browse)
        file_row.addWidget(self.path_edit)
        file_row.addWidget(browse)
        layout.addLayout(file_row)

        # Optional instruction
        layout.addWidget(QtWidgets.QLabel("Instruction (optional):"))
        self.instr_edit = QtWidgets.QTextEdit()
        self.instr_edit.setPlaceholderText("E.g. Focus on cardio only; make bullets tighter; prefer NICE over ACC…")
        self.instr_edit.setFixedHeight(100)
        layout.addWidget(self.instr_edit)

        # Action checkboxes
        layout.addWidget(QtWidgets.QLabel("Actions:"))
        self.cb_summary    = QtWidgets.QCheckBox("Create summary")
        self.cb_flashcards = QtWidgets.QCheckBox("Generate flashcards")
        self.cb_quiz       = QtWidgets.QCheckBox("Run quick quiz (5)")
        for cb in (self.cb_summary, self.cb_flashcards, self.cb_quiz):
            cb.setChecked(True)
            layout.addWidget(cb)

        # Buttons
        buttons = QtWidgets.QDialogButtonBox(QtWidgets.QDialogButtonBox.StandardButton.Ok |
                                             QtWidgets.QDialogButtonBox.StandardButton.Cancel)
        buttons.accepted.connect(self._on_ok)
        buttons.rejected.connect(self.reject)
        layout.addStretch(1)
        layout.addWidget(buttons)

    def _browse(self):
        path, _ = QtWidgets.QFileDialog.getOpenFileName(
            self, "Select file", "", "PowerPoint (*.pptx);;PDF (*.pdf)"
        )
        if path:
            self.path_edit.setText(path)

    def _on_ok(self):
        if not self.path_edit.text().strip():
            QtWidgets.QMessageBox.warning(self, "Missing file", "Please select a .pptx or .pdf file.")
            return
        # If user unticks everything, default to all 3
        if not (self.cb_summary.isChecked() or self.cb_flashcards.isChecked() or self.cb_quiz.isChecked()):
            self.cb_summary.setChecked(True)
            self.cb_flashcards.setChecked(True)
            self.cb_quiz.setChecked(True)
        self.accept()

    def values(self):
        return (
            self.path_edit.text().strip(),
            self.instr_edit.toPlainText().strip(),
            self.cb_summary.isChecked(),
            self.cb_flashcards.isChecked(),
            self.cb_quiz.isChecked(),
        )

    
class MyWindow(QMainWindow):

    def __init__(self, user_data):
        super(MyWindow, self).__init__()
        self.keyboard_enabled = True
        # store the current user data in the window
        self.user_data = user_data

        self.flashcards = []
        self.anki_filename = None
        self.conversation_history = []
        self.lecture_active = False
        self.lecture_pieces = []   # list of dicts: { "mode": "Explain/Summarise", "text": str, "ts": float }
        
        # Session storage for slides viewer
        self.saved_lectures = {}  # {lecture_name: {summaries: [], flashcards: [], created: timestamp}}
        self.load_saved_lectures()

        # Window settings
        self.setGeometry(100, 100, 325, 600)  # Width = 325 (narrow), Height = 600
        self.setMinimumSize(325, 550)         # Minimum allowed resize
        self.setMaximumWidth(325)             # Prevent horizontal expansion

        self.setWindowTitle('akson')
        self.setWindowFlags(
            
            QtCore.Qt.WindowType.WindowStaysOnTopHint
        )
        self.setWindowFlag(QtCore.Qt.WindowType.WindowDoesNotAcceptFocus, False)
        
        # Set app icon (macOS)
        self._setup_app_icon()
        
        # Create macOS menu bar
        self._create_menu_bar()
        
        # Setup macOS keyboard shortcuts
        self._setup_keyboard_shortcuts()
        
        # Detect macOS system theme and sync with webview
        self._setup_macos_theme_detection()
        
        
        # Apply exact site stylex
        self.setStyleSheet("""
            QWidget, QMainWindow{
                background-color: #1e1e1e;
                border-radius: 8px;
                font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            }
            
            QPushButton {
                background: transparent;
                color: #dcdcdc;
                padding: 6px 12px;
                border-radius: 8px;
                font-size: 16px;
                font-weight: 500;
                border: 0.5px solid #2a2a2e;
            }
            QPushButton:hover {
                
                color: #f5f5f5;
              
                font-size: 16px;
                border: 0.5px solid #f5f5f5;
            }
      

            QLabel {
                font-size: 20px;
                font-weight: 700;
                color: #dcdcdc;
                padding: 10px 6px;
            }
            QTextEdit {
                font-size: 14px;
                padding: 4px;
                border: none;
                border-radius: 8px;
                background-color: #1e1e1e;
                color: #dcdcdc;
                font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            }
        """)

        self.initUI()

                # Compact mini-dock (hidden by default)
        self.compact_bar = CompactBar(self)
        

        self.flashcards = []  
        self.keyboard_thread = threading.Thread(target=self.start_keyboard_listener, daemon=True)
        self.keyboard_thread.start()
        # Make window slightly transparent
        self.setWindowOpacity(0.98)



      
        QApplication.clipboard().clear()
        self._last_clip = ""

        self.clip_timer = QTimer(self)
        self.clip_timer.setInterval(300)                      
        self.clip_timer.timeout.connect(self.check_clipboard)
        self.clip_timer.start()
        print("[DEBUG] clipboard polling started")

          
        self.clipboard = QApplication.clipboard()
        
        # Check screen recording permissions on macOS after a short delay
        if sys.platform == "darwin":
            QtCore.QTimer.singleShot(1000, self.check_screen_permissions)
    @QtCore.pyqtSlot()


    @QtCore.pyqtSlot(str)
    def process_selected_text(self, text):
        # show loading
        self.loading_label.setText("Processing selection…")
        self.loading_label.show()
        # 1) Clear any old output so it's not mistaken for new
        self.text_box.clear()
        # 2) Force Qt to repaint the cleared box immediately
        QtWidgets.QApplication.processEvents()
        # 3) Now show loading
        self.loading_label.setText("Processing selection…")
        self.loading_label.show()
        self.last_extracted_text = text


        def worker():
            try:
                mode = self.mode_selector.currentText()
                if mode == "Explain":
                    prompt = f"""You are Akson, a study co-pilot and educator. Your job is to convert a raw lecture slide into a high-yield, deeply explained, practical, exam-ready teaching resource. You're not summarizing — you're teaching.

🧠 CORE INSTRUCTIONS

If only one or two or a couple of words are given, then DEFINE the words, even if multiple words are given, start by shortly defining the words and where they come from. THIS IS CRUCIAL!!! After that give a line break and continue;
1️⃣ Read & Filter Smartly
Carefully read the entire slide. Completely ignore accidental or irrelevant text. Focus only on important content.

2️⃣ Teach, Don't Rephrase
Do not paraphrase or simplify blindly. Instead, explain everything clearly and intelligently — as if teaching a bright but uncertain student. For every concept, mechanism, process, term, abbreviation, or instruction:

Define it (briefly and clearly)
Explain what it does
Explain why it matters in practice
Add what the slide leaves out but a good teacher would include
Assume the student has knowledge gaps. Fill them proactively.

3️⃣ Fill In What the Slide Leaves Out
If the slide vaguely says things like "evaluation," "steps," or "solutions," you must:

List what those tests might be
Briefly explain why they're ordered
Include standard options or typical approaches where appropriate
Never allow vague or incomplete statements to pass through unexplained.
4️⃣ Always Add Real-World Context
Every explanation must tie into real-world applications. Show how the content connects to:

Key signals or indicators
Reasoning and decision steps
Investigations or tools
Solutions or approaches (baseline vs. advanced)
Outcomes
Mechanisms and potential pitfalls
5️⃣ Prioritise High-Yield Content
Mark 🔑 high-yield facts, 💡 mnemonics or tips, and ⚠️ exam traps or common misconceptions.
Mention anything that's commonly tested or clinically critical. Highlight contrasts and "commonly confused" points.

6️⃣ Use Active Recall Sparingly
After each major idea or section, insert 1 short active recall question (if appropriate).
Example: 💭 "What's the main risk factor for this condition?"
Keep questions simple and relevant.

7️⃣ Structure for Maximum Learning
Use a clean, consistent layout:

Clear section headings
Bullet points
Line spacing for readability
Concise but precise language
Avoid waffle. Use emojis (🧠, 💉, ⚠️, ❤️) sparingly to boost engagement.
8️⃣ Always Add Value Beyond the Slide
If the slide is shallow or incomplete, you must enrich it. Fill in mechanisms. Suggest simple mnemonics. Clarify unclear pathways. Break complex concepts into digestible steps. Think: "If the student sees this in a viva or ward round, what would they need to say or understand?"

9️⃣ Your Mental Model
You are a sharp, exam-aware tutor.
You are here to build true understanding, not copy text.
The final output must help the student:

Grasp the concept
Retain it
Apply it
Pass their exams with confidence

If TEXT IS SELECTED or hovered over or highlighted:
only explain the word selected, or the phrase or sentence(s) selected.


        {text.strip()}

        """
                else:
                    prompt = f"""You summarise OCR text from lecture slides into SHORT, high-yield bullets for exam revision.

GOAL
Return a compact list of independent facts in arrow/abbreviation style (memory notes), not prose.

HARD RULES (must obey)
 - start with a heading so we know what you will speak about.
 - Make it short but easy to understand. use words and wording that makes it easy to comprehend

- 4–6 bullets total. One fact per bullet. No paragraphs.
- Do NOT start bullets with category labels (e.g., "Epidemiology:", "Prevalence:", "Demographics:", "Skeletal destruction:").
  • Allowed labels only when logically required: **Criteria** and **Comparison**.
- Never chain multiple facts with ";" or long commas. Split into separate bullets.
- Use symbols & shorthand aggressively:
  ↑ increased/elevated, ↓ decreased/low, → leads to/causes, ~ approx, ≈ ratio, κ/λ light chains.
  - Prefer standard abbreviations when obvious from context.
- Ratios compact: "M:F ≈ 3:1", "Blacks:Whites ≈ 2:1".
- Minimal bold: only the key term in a bullet (e.g., **NSTEMI**, **Troponin I**, **Russell bodies**).
- Zero filler or basic definitions.

REASONING (internal; do not output)
1) Identify intent (definition/essence, criteria, comparison, atypical features, pathology/morphology, epidemiology).
2) Select the 4–6 highest-yield, exam-testable facts; prefer differentiators and criteria.
3) Rewrite each as: **Label** → compressed fact(s) with symbols/abbrevs.
4) Self-check before output:
   - Bullet count 4–6 or as little as possible.
   - No bullet starts with a category label (except **Criteria**/**Comparison** when necessary).
   - No ";" inside bullets.
   - Each bullet ≤ ~18–22 words.

FORMATS WHEN NEEDED
- **Criteria:** Use a single bullet if it's the classic triad/"≥2 of 3". Otherwise split into bullets.
- **Comparison:** Up to 3 bullets, one per entity, each starting with a bold entity then its defining line.

Return ONLY the bullets.


{text.strip()}
"""
    



                # === Append custom prompt if user defined one ===
                custom_prompt = self.user_data.get("custom_prompt", "")
                system_msg = "You are a helpful study tutor."
                if custom_prompt:
                    system_msg += (
                        "\n\n# USER OVERRIDE INSTRUCTIONS (highest priority):\n"
                        f"{custom_prompt}\n"
                        "If any instruction above conflicts with previous instructions, "
                        "FOLLOW THE USER OVERRIDE INSTRUCTIONS."
                    )

                response = openai.ChatCompletion.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": system_msg},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.7,
                    stream=True
                )


                summary = ""
                
                for chunk in response:
                    delta = chunk["choices"][0]["delta"].get("content", "")
                    summary += delta

                    # optional: live update every N chars
                    if len(summary) % 100 == 0:
                        html = markdown_to_html(summary)
                        QtCore.QMetaObject.invokeMethod(
                            self.text_box, "setHtml",
                            QtCore.Qt.ConnectionType.QueuedConnection,
                            QtCore.Q_ARG(str, html)
                        )
                # right after finalizing 'summary'
                self.latest_context = summary.strip()
                # collect for lecture pack



                # final render
                html = markdown_to_html(summary.strip())
                QtCore.QMetaObject.invokeMethod(
                    self.text_box, "setHtml",
                    QtCore.Qt.ConnectionType.QueuedConnection,
                    QtCore.Q_ARG(str, html)

                )
                QtCore.QMetaObject.invokeMethod(
                    self, "_collect_lecture_piece",
                    QtCore.Qt.ConnectionType.QueuedConnection,
                    QtCore.Q_ARG(str, summary.strip()),
                    QtCore.Q_ARG(str, mode)
                )

            except Exception as e:
                QtCore.QMetaObject.invokeMethod(
                    self.text_box, "setPlainText",
                    QtCore.Qt.ConnectionType.QueuedConnection,
                    QtCore.Q_ARG(str, f"Error: {e}")
                )
            finally:
                QtCore.QMetaObject.invokeMethod(
                    self.loading_label, "hide",
                    QtCore.Qt.ConnectionType.QueuedConnection
                )

        threading.Thread(target=worker, daemon=True).start()


    @pyqtSlot()
    @pyqtSlot()


    @pyqtSlot()
    @pyqtSlot()
    def check_clipboard(self):
        text = QApplication.clipboard().text().strip()
        if not text or text == self._last_clip:
            return
        self._last_clip = text
        print(f"[DEBUG] new clipboard text → {repr(text)})")

        if hasattr(self, "compact_bar") and self.compact_bar.isVisible():
            self.compact_bar.explain_clipboard(text)
        else:
            self.process_selected_text(text)



    def initUI(self):
        main_layout = QtWidgets.QVBoxLayout()
        main_layout.setContentsMargins(8, 8, 8, 8)
        main_layout.setSpacing(8)


        self.latest_context = ""
        self.conversation_history = []

        # Mode Selector
        mode_layout = QtWidgets.QHBoxLayout()
        mode_label = QtWidgets.QLabel("")
        mode_label.setStyleSheet("font-size: 12px; color: #8e8e93; padding: 0 4px 0 0;")

        self.mode_selector = QtWidgets.QComboBox()
        self.mode_selector.addItems(["Explain", "Summary"])
        
        self.mode_selector.setFixedWidth(90)  # Reduced from 110
        self.mode_selector.setStyleSheet("""
                background: transparent;
                color: #dcdcdc;
                border: none;
                font-size: 12px;
                padding: 2px 4px;
           
        """)
        #self.label = QtWidgets.QLabel("")
        #self.label.setStyleSheet("color: gray; font-size: 13px; padding-left: 6px;")
        #main_layout.addWidget(self.label)

    
    
        mode_layout.addWidget(self.mode_selector)
        mode_layout.addStretch()

        self.import_slides_button = QtWidgets.QPushButton('↑')
        # Settings Button (right side)
        self.settings_button = QtWidgets.QPushButton("⚙")

        self.settings_button.setCursor(QtCore.Qt.CursorShape.PointingHandCursor)
        self.settings_button.setFixedWidth(28)
        
        self.settings_button.setStyleSheet("""
            QPushButton {
                background: transparent;
                color: #dcdcdc;
                border: none;
                font-size: 15px;
                padding: 2px;
            }
        """)
        self.settings_button.clicked.connect(self.open_settings)
        self.import_slides_button.clicked.connect(self.open_import_dialog)
        
        # Add slides viewer button
        self.slides_viewer_button = QtWidgets.QPushButton('Slides')
        self.slides_viewer_button.setToolTip("Document Viewer")
        self.slides_viewer_button.clicked.connect(self.switch_to_slides_tab)
        self.slides_viewer_button.setCheckable(True)
        self.slides_viewer_button.setStyleSheet("""
            QPushButton {
                background: transparent;
                color: #dcdcdc;
                border: 1px solid #333;
                border-radius: 4px;
                padding: 4px 6px;
                font-size: 11px;
                font-weight: 600;
            }
            QPushButton:hover {
                background: #2a2a2a;
                border-color: #555;
            }
            QPushButton:checked {
                background: #3a3a3a;
                border-color: #666;
            }
        """)
        
        # Add main app button
        self.main_app_button = QtWidgets.QPushButton('Home')
        self.main_app_button.setToolTip("Main App")
        self.main_app_button.clicked.connect(self.switch_to_main_tab)
        self.main_app_button.setCheckable(True)
        self.main_app_button.setChecked(True)
        self.main_app_button.setStyleSheet("""
            QPushButton {
                background: transparent;
                color: #dcdcdc;
                border: 1px solid #333;
                border-radius: 4px;
                padding: 4px 6px;
                font-size: 11px;
                font-weight: 600;
            }
            QPushButton:hover {
                background: #2a2a2a;
                border-color: #555;
            }
            QPushButton:checked {
                background: #3a3a3a;
                border-color: #666;
            }
        """)
        
                # Mini-dock trigger (additional minimize)
        self.mini_button = QtWidgets.QPushButton("–")
        self.mini_button.setCursor(QtCore.Qt.CursorShape.PointingHandCursor)
        self.mini_button.setToolTip("Mini-dock")
        self.mini_button.setFixedWidth(28)
        self.mini_button.setStyleSheet("""
            QPushButton {
                background: transparent;
                color: #dcdcdc;
                border: none;
                font-size: 15px;
                padding: 2px;
            }
            QPushButton:hover {
                color: #d0caff;
            }
        """)
        self.import_slides_button.setFixedWidth(28)
        self.import_slides_button.setStyleSheet("""
            QPushButton {
                background: transparent;
                color: #dcdcdc;
                border: none;
                font-size: 15px;
                padding: 2px;
            }
            QPushButton:hover {
                color: #d0caff;
            }
        """)
        

        


        self.mini_button.clicked.connect(self.enter_compact_mode)
        mode_layout.addWidget(self.main_app_button)
        mode_layout.addWidget(self.import_slides_button)
        mode_layout.addWidget(self.slides_viewer_button)
        mode_layout.addWidget(self.settings_button)
        mode_layout.addWidget(self.mini_button)
  


        main_layout.addLayout(mode_layout)

         
        # === Summarize button ===
        # Create grid layout for buttons
       
        self.start_lecture_button = QtWidgets.QPushButton('Start')
        self.create_flashcard_button = QtWidgets.QPushButton('Flashcard')
        self.end_lecture_button = QtWidgets.QPushButton('End')
        

        self.end_lecture_button.setEnabled(False)
  
        


        top_buttons_layout = QtWidgets.QHBoxLayout()
        top_buttons_layout.setSpacing(5)

        # Make buttons expand and divide space equally
        for btn in [self.start_lecture_button, self.create_flashcard_button, self.end_lecture_button]:
    # (styling stays exactly as you already have it)

            btn.setSizePolicy(QtWidgets.QSizePolicy.Policy.Preferred, QtWidgets.QSizePolicy.Policy.Fixed)
            btn.setMinimumWidth(0)              # allow shrinking
            btn.setFixedHeight(26)
            
            btn.setStyleSheet("""
                QPushButton {
                    background: transparent;
                    color: #dcdcdc;
                    font-size: 12px;
                    font-weight: 600;
                    border: 0.5px solid #2a2a2e;
                    border-radius: 8px;
                    padding: 2px 6px;
                    min-width: 0;
                }
                QPushButton:hover {
                    background-color: #1e1e1e;
                    border: 0.5px solid #7f7fff;
                    color: #d0caff;
                }
            """)

            top_buttons_layout.addWidget(btn)



        # Connect button signals
        #self.summarize_button.clicked.connect(self.manual_summarize)
        #self.clear_button.clicked.connect(self.clear_summary)
        #self.ask_button.clicked.connect(self.ask_question)
        self.start_lecture_button.clicked.connect(self.start_lecture)
        self.create_flashcard_button.clicked.connect(self.create_flashcard)
        self.end_lecture_button.clicked.connect(self.end_lecture)
        


        # Smaller button sizes & style
        
        main_layout.addLayout(top_buttons_layout)

        # Grid positions
 

        # Disable flashcard buttons at start
        self.create_flashcard_button.setEnabled(False)

        # Add grid to main layout



        # === Output box ===
        # === Composite Summary Frame (Text + Buttons) ===
        summary_frame = QtWidgets.QFrame()
        summary_frame.setStyleSheet("""
            QFrame {
                background-color: #1e1e1e;
                border-radius: 8px;
            }
        """)
        summary_layout = QtWidgets.QVBoxLayout(summary_frame)
        summary_layout.setContentsMargins(8, 8, 8, 8)
        summary_layout.setSpacing(8)

        # Summary text box
        self.text_box = QtWidgets.QTextBrowser()
        self.text_box.setVerticalScrollBarPolicy(QtCore.Qt.ScrollBarPolicy.ScrollBarAlwaysOff)
        self.text_box.setHorizontalScrollBarPolicy(QtCore.Qt.ScrollBarPolicy.ScrollBarAlwaysOff)
        self.text_box.setOpenExternalLinks(True)
        self.text_box.setMinimumHeight(335)
        self.text_box.setMaximumWidth(309)  # Ensure content doesn't overflow (325 - 16 for margins)
        self.text_box.setWordWrapMode(QtGui.QTextOption.WrapMode.WrapAtWordBoundaryOrAnywhere)
        self.text_box.setHtml("<span style='color: gray;'>Content will appear here...</span>")
        self.text_box.setStyleSheet("""
            QTextBrowser {
                font-size: 13.5px;
                padding: 4px;
                border: none;
                background-color: transparent;
                color: #dcdcdc;
                font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            }
        """)
        summary_layout.addWidget(self.text_box)

        # Buttons inside the summary box
        summary_button_row = QtWidgets.QHBoxLayout()
        summary_button_row.setSpacing(6)

        
        self.clear_button = QtWidgets.QPushButton("Clear")
        self.summarize_button = QtWidgets.QPushButton("Go")
        


 
       # ───────────────────────────

        

        for btn in [self.summarize_button, self.clear_button]:
            btn.setCursor(QtCore.Qt.CursorShape.PointingHandCursor)
            btn.setFixedHeight(26)
            btn.setStyleSheet("""
                QPushButton {
                    background: transparent;
                    color: #dcdcdc;
                    font-size: 12.5px;
                    font-weight: 600;
                    border: 0.5px solid #2a2a2e;
                    border-radius: 8px;
                    
                }
                QPushButton:hover {
                    background-color: #1e1e1e;
                    border: 0.5px solid #7f7fff;
                    color: #d0caff;
                }
            """)
            summary_button_row.addWidget(btn)

        summary_layout.addLayout(summary_button_row)

        # Add to main layout
        main_layout.addWidget(summary_frame)


        # Move the existing Ask button here, below the summary box
        #main_layout.addWidget(self.ask_button)




                # === Action Buttons under summary ===
       
        self.summarize_button.clicked.connect(self.manual_summarize)
        self.clear_button.clicked.connect(self.clear_summary)
        #self.ask_button.clicked.connect(self.ask_question)
        self.loading_label = QtWidgets.QLabel("Processing...")
        self.loading_label.setAlignment(QtCore.Qt.AlignmentFlag.AlignCenter)
        self.loading_label.setStyleSheet("color: #888; font-size: 13px; padding: 4px;")
        self.loading_label.hide()
        


                # ─── Inline Ask Input ───
        # ─── Inline Ask Input with Arrow ───
        self.ask_input = QtWidgets.QLineEdit()
        self.ask_input.setPlaceholderText("Ask anything…")
        self.ask_input.setMaximumWidth(309)  # Constrain width
        self.ask_input.setStyleSheet("""
            QLineEdit {
                background-color: #1e1e1e;
                color: #dcdcdc;
                border: 0.5px solid #2a2a2e;
                border-radius: 8px;
                padding: 6px 32px 6px 6px;
                font-size: 12px;
            }
        """)
        # Make it taller than a default single line
        self.ask_input.setFixedHeight(34)

        self.ask_input.setSizePolicy(
            QtWidgets.QSizePolicy.Policy.Expanding,
            QtWidgets.QSizePolicy.Policy.Fixed
        )
        self.ask_input.setFocusPolicy(QtCore.Qt.FocusPolicy.StrongFocus)

        # Add trailing arrow action
        arrow_icon = QtGui.QIcon.fromTheme("go-next")  # or QtGui.QIcon("path/to/your/arrow.png")
        send_action = self.ask_input.addAction(
            arrow_icon,
            QtWidgets.QLineEdit.ActionPosition.TrailingPosition
        )
        send_action.triggered.connect(lambda: self._inline_ask())

        # Also send on Enter
        self.ask_input.returnPressed.connect(lambda: self._inline_ask())

        main_layout.addWidget(self.ask_input)
        main_layout.addWidget(self.loading_label)
# ─────────────────────────────────────

        # ───────────────────────────



        scroll_hint = QtWidgets.QLabel("⬇ Scroll down for flashcards ⬇")
        scroll_hint.setAlignment(QtCore.Qt.AlignmentFlag.AlignCenter)
        scroll_hint.setStyleSheet("color: gray; font-size: 13px; padding: 6px;")
        main_layout.addWidget(scroll_hint)


        flashcard_title = QtWidgets.QLabel("Flashcards:")
        flashcard_title.setStyleSheet("color: #dcdcdc; font-size: 16px; font-weight: bold; padding: 4px;")
        
        main_layout.addWidget(flashcard_title)


        # Flashcard review list
        self.flashcard_list = QtWidgets.QListWidget()

        self.flashcard_list.setWordWrap(True)
        self.flashcard_list.setMinimumHeight(400)


        self.flashcard_list.setSizePolicy(QtWidgets.QSizePolicy.Policy.Expanding, QtWidgets.QSizePolicy.Policy.Maximum)
        
        self.flashcard_list.setSizePolicy(QtWidgets.QSizePolicy.Policy.Expanding, QtWidgets.QSizePolicy.Policy.Expanding)
        main_layout.addWidget(self.flashcard_list, stretch=1)

        # Library section
        library_title = QtWidgets.QLabel("Saved Lectures:")
        library_title.setStyleSheet("color: #dcdcdc; font-size: 16px; font-weight: bold; padding: 4px; margin-top: 10px;")
        main_layout.addWidget(library_title)

        self.library_list = QtWidgets.QListWidget()
        self.library_list.setWordWrap(True)
        self.library_list.setMinimumHeight(200)
        self.library_list.setStyleSheet("""
            QListWidget {
                background-color: #2a2a2d;
                border: 1px solid #333;
                border-radius: 6px;
                padding: 8px;
                color: #ececec;
                font-size: 13px;
            }
            QListWidget::item {
                padding: 8px;
                border-bottom: 1px solid #333;
            }
r            QListWidget::item:selected {
                background-color: #3a3a3d;
            }
        """)
        self.library_list.setSizePolicy(QtWidgets.QSizePolicy.Policy.Expanding, QtWidgets.QSizePolicy.Policy.Expanding)
        main_layout.addWidget(self.library_list, stretch=1)

        # Library buttons
        library_buttons = QtWidgets.QHBoxLayout()
        self.view_lecture_btn = QtWidgets.QPushButton("View Lecture")
        self.delete_lecture_btn = QtWidgets.QPushButton("Delete")
        
        self.view_lecture_btn.setStyleSheet("""
            QPushButton {
                background: #3a3a3d;
                border: 1px solid #555;
                color: #ececec;
                padding: 6px 12px;
                border-radius: 4px;
                font-size: 12px;
            }
            QPushButton:hover {
                background: #4a4a4d;
            }
        """)
        
        self.delete_lecture_btn.setStyleSheet("""
            QPushButton {
                background: #d32f2f;
                border: 1px solid #f44336;
                color: white;
                padding: 6px 12px;
                border-radius: 4px;
                font-size: 12px;
            }
            QPushButton:hover {
                background: #f44336;
            }
        """)
        
        self.view_lecture_btn.clicked.connect(self.view_selected_lecture)
        self.delete_lecture_btn.clicked.connect(self.delete_selected_lecture)
        
        library_buttons.addWidget(self.view_lecture_btn)
        library_buttons.addWidget(self.delete_lecture_btn)
        library_buttons.addStretch()
        
        main_layout.addLayout(library_buttons)


        # Wrap everything into a scrollable container
        # Wrap everything into a scrollable container
        container = QtWidgets.QWidget()
        container.setLayout(main_layout)
        container.setSizePolicy(QtWidgets.QSizePolicy.Policy.MinimumExpanding, QtWidgets.QSizePolicy.Policy.MinimumExpanding)
        container.setMinimumWidth(0)
        container.setMaximumWidth(325)  # Prevent horizontal expansion
        container.setMinimumHeight(0)


        scroll = QtWidgets.QScrollArea()
        scroll.setWidgetResizable(True)
        scroll.setSizePolicy(QtWidgets.QSizePolicy.Policy.Expanding, QtWidgets.QSizePolicy.Policy.Expanding)
        scroll.setVerticalScrollBarPolicy(QtCore.Qt.ScrollBarPolicy.ScrollBarAsNeeded)
        scroll.setHorizontalScrollBarPolicy(QtCore.Qt.ScrollBarPolicy.ScrollBarAlwaysOff)

        scroll.setWidget(container)
        scroll.setStyleSheet("""
    QScrollBar:vertical, QScrollBar:horizontal {
        width: 0px;
        height: 0px;
        background: transparent;
    }
""")


        self.setCentralWidget(scroll)
        
        # Initialize library display
        self.update_library_display()


    def open_import_dialog(self):
        dlg = ImportFilesDialog(self)
        if dlg.exec():
            path, instr, do_summary, do_flashcards, do_quiz = dlg.values()
            self._start_import_job(path, instr, do_summary, do_flashcards, do_quiz)

    def switch_to_slides_tab(self):
        """Switch to slides/document viewer tab"""
        try:
            if not hasattr(self, 'slides_webview'):
                self.create_slides_webview()
            
            # Hide main content and show slides
            self.text_box.hide()
            self.flashcard_list.hide()
            self.slides_webview.show()
            
            # Update button states
            self.main_app_button.setChecked(False)
            self.slides_viewer_button.setChecked(True)
            
        except Exception as e:
            print(f"Error switching to slides tab: {e}")
            self.text_box.setPlainText(f"Error switching to document viewer: {e}")

    def switch_to_main_tab(self):
        """Switch to main app tab"""
        try:
            # Show main content and hide slides
            self.text_box.show()
            self.flashcard_list.show()
            if hasattr(self, 'slides_webview'):
                self.slides_webview.hide()
            
            # Update button states
            self.main_app_button.setChecked(True)
            self.slides_viewer_button.setChecked(False)
            
        except Exception as e:
            print(f"Error switching to main tab: {e}")

    def create_slides_webview(self):
        """Create the slides webview widget"""
        try:
            self.slides_webview = QWebEngineView()
            self.slides_webview.setStyleSheet("""
                QWebEngineView {
                    background-color: #1a1a1d;
                    border: none;
                }
            """)
            
            # Load the slides HTML content
            slides_html = self.get_slides_html()
            self.slides_webview.setHtml(slides_html)
            
            # Add JavaScript bridge for tab switching and session saving
            self.slides_webview.page().runJavaScript("""
                window.switchToMainTab = function() {
                    // This will be handled by the parent
                    console.log('Switching to main tab');
                };
                
                window.saveSession = function(lectureName, summaries, flashcards) {
                    // This will be handled by the parent
                    console.log('Saving session:', lectureName, summaries.length, flashcards.length);
                };
            """)
            
            # Connect the JavaScript bridge
            self.slides_webview.page().javaScriptConsoleMessage = self.handle_js_console_message
            
            # Add to main layout
            self.main_layout.addWidget(self.slides_webview)
            self.slides_webview.hide()  # Initially hidden
            
        except Exception as e:
            print(f"Error creating slides webview: {e}")

    def handle_js_console_message(self, level, message, lineNumber, sourceID):
        """Handle JavaScript console messages for tab switching and session saving"""
        if "Switching to main tab" in message:
            self.switch_to_main_tab()
        elif "Saving session:" in message:
            # Extract session data from the webview
            self.slides_webview.page().runJavaScript("""
                var sessionData = {
                    lectureName: window.pendingLectureName || 'Untitled Lecture',
                    summaries: window.collectedSummaries || [],
                    flashcards: window.collectedFlashcards || []
                };
                console.log('SESSION_DATA:' + JSON.stringify(sessionData));
            """)
        elif "SESSION_DATA:" in message:
            # Handle the session data
            try:
                import json
                data_str = message.split('SESSION_DATA:')[1]
                session_data = json.loads(data_str)
                
                # Get lecture name from the data
                lecture_name = session_data.get('lectureName', 'Untitled Lecture')
                
                # Save the session data
                success = self.save_session_data(lecture_name, session_data.get('summaries', []), session_data.get('flashcards', []))
                
                if success:
                    self.text_box.setHtml(f"<span style='color:#4caf50'>✅ Lecture '{lecture_name}' saved successfully!</span>")
                else:
                    self.text_box.setHtml(f"<span style='color:#f44336'>❌ Error saving lecture '{lecture_name}'</span>")
                    
            except Exception as e:
                print(f"Error handling session data: {e}")

    def load_saved_lectures(self):
        """Load saved lectures from file"""
        try:
            import json
            import os
            lectures_file = os.path.join(os.path.dirname(__file__), "saved_lectures.json")
            if os.path.exists(lectures_file):
                with open(lectures_file, 'r', encoding='utf-8') as f:
                    self.saved_lectures = json.load(f)
        except Exception as e:
            print(f"Error loading saved lectures: {e}")
            self.saved_lectures = {}

    def save_saved_lectures(self):
        """Save lectures to file"""
        try:
            import json
            import os
            lectures_file = os.path.join(os.path.dirname(__file__), "saved_lectures.json")
            with open(lectures_file, 'w', encoding='utf-8') as f:
                json.dump(self.saved_lectures, f, indent=2, ensure_ascii=False)
        except Exception as e:
            print(f"Error saving lectures: {e}")

    def save_session_data(self, lecture_name, summaries, flashcards):
        """Save session data (summaries and flashcards) for a lecture"""
        try:
            import time
            from datetime import datetime
            
            # Deduplicate summaries
            unique_summaries = self.deduplicate_summaries(summaries)
            
            # Deduplicate flashcards
            unique_flashcards = self.deduplicate_flashcards(flashcards)
            
            # Save to lectures
            self.saved_lectures[lecture_name] = {
                'summaries': unique_summaries,
                'flashcards': unique_flashcards,
                'created': datetime.now().isoformat(),
                'updated': datetime.now().isoformat()
            }
            
            # Save to file
            self.save_saved_lectures()
            
            # Update UI
            self.update_library_display()
            
            return True
        except Exception as e:
            print(f"Error saving session data: {e}")
            return False

    def deduplicate_summaries(self, summaries):
        """Remove duplicate and obsolete summaries"""
        if not summaries:
            return []
        
        unique_summaries = []
        seen_content = set()
        
        for summary in summaries:
            # Normalize content for comparison
            content = summary.get('content', '').strip().lower()
            if content and content not in seen_content:
                seen_content.add(content)
                unique_summaries.append(summary)
        
        return unique_summaries

    def deduplicate_flashcards(self, flashcards):
        """Remove duplicate flashcards"""
        if not flashcards:
            return []
        
        unique_flashcards = []
        seen_cards = set()
        
        for card in flashcards:
            # Create a key based on question and answer
            key = f"{card.get('question', '').strip().lower()}|{card.get('answer', '').strip().lower()}"
            if key not in seen_cards:
                seen_cards.add(key)
                unique_flashcards.append(card)
        
        return unique_flashcards

    def update_library_display(self):
        """Update the library display in the main app"""
        try:
            if hasattr(self, 'library_list'):
                self.library_list.clear()
                for lecture_name, data in self.saved_lectures.items():
                    summary_count = len(data.get('summaries', []))
                    flashcard_count = len(data.get('flashcards', []))
                    created = data.get('created', 'Unknown')
                    
                    item_text = f"{lecture_name}\n{summary_count} summaries | {flashcard_count} flashcards\nCreated: {created[:10]}"
                    self.library_list.addItem(item_text)
        except Exception as e:
            print(f"Error updating library display: {e}")

    def view_selected_lecture(self):
        """View the selected lecture's summaries and flashcards"""
        try:
            current_item = self.library_list.currentItem()
            if not current_item:
                return
            
            # Extract lecture name from the item text
            lecture_name = current_item.text().split('\n')[0]
            
            if lecture_name in self.saved_lectures:
                data = self.saved_lectures[lecture_name]
                summaries = data.get('summaries', [])
                flashcards = data.get('flashcards', [])
                
                # Display summaries
                summary_text = f"{lecture_name}\n\n"
                summary_text += "SUMMARIES:\n" + "="*50 + "\n\n"
                
                for i, summary in enumerate(summaries, 1):
                    summary_text += f"{i}. {summary.get('content', '')}\n\n"
                
                summary_text += "\nFLASHCARDS:\n" + "="*50 + "\n\n"
                
                for i, card in enumerate(flashcards, 1):
                    summary_text += f"{i}. Q: {card.get('question', '')}\n"
                    summary_text += f"   A: {card.get('answer', '')}\n\n"
                
                # Display in text box
                self.text_box.setHtml(markdown_to_html(summary_text))
                
        except Exception as e:
            print(f"Error viewing lecture: {e}")

    def delete_selected_lecture(self):
        """Delete the selected lecture"""
        try:
            current_item = self.library_list.currentItem()
            if not current_item:
                return
            
            # Extract lecture name
            lecture_name = current_item.text().split('\n')[0]
            
            # Confirm deletion
            from PyQt6.QtWidgets import QMessageBox
            reply = QMessageBox.question(
                self, 
                'Delete Lecture', 
                f'Are you sure you want to delete "{lecture_name}"?',
                QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No,
                QMessageBox.StandardButton.No
            )
            
            if reply == QMessageBox.StandardButton.Yes:
                if lecture_name in self.saved_lectures:
                    del self.saved_lectures[lecture_name]
                    self.save_saved_lectures()
                    self.update_library_display()
                    
        except Exception as e:
            print(f"Error deleting lecture: {e}")

    def get_slides_html(self):
        """Get the HTML content for the slides viewer"""
        # This will be the slides.py HTML content integrated here
        return """
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="utf-8">
            <title>Document Viewer</title>
            <style>
                body { 
                    margin: 0; 
                    padding: 0; 
                    background: #1a1a1d; 
                    color: #ececee; 
                    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
                }
                #root { 
                    display: grid; 
                    grid-template-columns: 1fr 300px; 
                    height: 100vh; 
                }
                #pdfFrame { 
                    width: 100%; 
                    height: 100%; 
                    border: none; 
                    background: #2a2a2d; 
                }
                #rightSidebar { 
                    background: #1a1a1d; 
                    border-left: 1px solid #333; 
                    padding: 16px; 
                    overflow-y: auto; 
                }
                .rsSection { 
                    margin-bottom: 20px; 
                }
                .rsSection h3 { 
                    color: #d8d8da; 
                    font-size: 14px; 
                    margin: 0 0 8px 0; 
                    font-weight: 600; 
                }
                .rsBox { 
                    background: #2a2a2d; 
                    border: 1px solid #333; 
                    border-radius: 6px; 
                    padding: 12px; 
                    font-size: 13px; 
                    line-height: 1.4; 
                    color: #ececee; 
                    white-space: pre-wrap; 
                }
                .rsEmpty { 
                    color: #666; 
                    font-style: italic; 
                }
                .bigBtn { 
                    background: #3a3a3d; 
                    border: 1px solid #555; 
                    color: #ececee; 
                    padding: 8px 12px; 
                    border-radius: 6px; 
                    cursor: pointer; 
                    font-size: 13px; 
                }
                .bigBtn:hover { 
                    background: #4a4a4d; 
                }
                #topBar { 
                    display: flex; 
                    gap: 10px; 
                    margin-bottom: 12px; 
                }
            </style>
        </head>
        <body>
            <div id="root">
                <iframe id="pdfFrame" src="about:blank"></iframe>
                <aside id="rightSidebar">
                    <div id="rsScroll">
                        <div id="topBar">
                            <button class="bigBtn" id="btnCompact" title="Back to Main App">Home</button>
                            <button class="bigBtn" id="btnSlides" title="Start presentation mode">Slides</button>
                            <button class="bigBtn" id="btnFlash" title="Open Akson Flashcards">Cards</button>
                            <button class="bigBtn" id="btnLibrary" title="View Library">Library</button>
                            <button class="bigBtn" id="btnEndSession" title="End Session & Save" style="background: #333; border-color: #555; color: #d8d8da;">Save</button>
                        </div>
                        
                        <section class="rsSection">
                            <h3>Current Page Summary</h3>
                            <div id="rsSummary" class="rsBox rsEmpty">Change page to generate a summary here.</div>
                        </section>

                        <section class="rsSection">
                            <h3>Term Explainer (select a single word)</h3>
                            <div id="rsExplain" class="rsBox rsEmpty">Select a term in the PDF to see a concise definition.</div>
                        </section>

                        <section class="rsSection">
                            <h3>Ask AI</h3>
                            <div style="display:flex;gap:6px;margin-bottom:8px;">
                                <input type="text" id="questionInput" placeholder="Ask a question..." style="flex:1;padding:6px;border-radius:6px;border:1px solid #333;background:#1a1a1d;color:#eee;font-size:13px;">
                                <button id="btnAsk" class="bigBtn" style="padding:6px 10px;">Ask</button>
                            </div>
                            <div id="aiAnswer" class="rsBox rsEmpty">Ask me anything about the document or general medical questions.</div>
                        </section>

                        <section class="rsSection">
                            <h3>Auto-Generated Flashcards</h3>
                            <div id="autoFcList"></div>
                            <div id="autoFcEmpty" class="rsEmpty">Flashcards will be generated automatically for relevant pages.</div>
                        </section>

                        <section class="rsSection">
                            <h3>Manual Flashcards</h3>
                            <div id="fcList"></div>
                            <div id="fcEmpty" class="rsEmpty">No cards yet.</div>
                            <div style="margin-top:8px;display:flex;gap:6px">
                                <button id="btnAddCard" class="bigBtn" style="padding:6px 10px;">+ Empty Card</button>
                                <button id="btnExportCards" class="bigBtn" style="padding:6px 10px;">Export (JSON)</button>
                            </div>
                        </section>
                    </div>
                </aside>
            </div>
            
            <script>
                // Session data collection
                window.collectedSummaries = [];
                window.collectedFlashcards = [];
                window.currentPage = 1;
                
                // Basic functionality for the slides viewer
                document.getElementById('btnCompact').onclick = function() {
                    // Call the parent's switch function
                    if (window.switchToMainTab) {
                        window.switchToMainTab();
                    }
                };
                
                document.getElementById('btnAsk').onclick = function() {
                    const question = document.getElementById('questionInput').value.trim();
                    if (question) {
                        document.getElementById('aiAnswer').textContent = 'Thinking...';
                        // This would be handled by the parent app's API
                    }
                };
                
                document.getElementById('btnEndSession').onclick = function() {
                    const lectureName = prompt('Enter lecture name to save:');
                    if (lectureName && lectureName.trim()) {
                        // Store lecture name for the parent
                        window.pendingLectureName = lectureName.trim();
                        
                        // Collect all summaries and flashcards
                        const summaries = window.collectedSummaries || [];
                        const flashcards = window.collectedFlashcards || [];
                        
                        // Send to parent app for saving
                        if (window.saveSession) {
                            window.saveSession(lectureName.trim(), summaries, flashcards);
                        }
                        
                        // Clear session data
                        window.collectedSummaries = [];
                        window.collectedFlashcards = [];
                    }
                };
                
                // Function to add summary to collection
                window.addSummary = function(page, content) {
                    if (content && content.trim()) {
                        window.collectedSummaries.push({
                            page: page,
                            content: content.trim(),
                            timestamp: new Date().toISOString()
                        });
                    }
                };
                
                // Function to add flashcards to collection
                window.addFlashcards = function(page, flashcards) {
                    if (flashcards && flashcards.length > 0) {
                        flashcards.forEach(card => {
                            window.collectedFlashcards.push({
                                page: page,
                                question: card.question,
                                answer: card.answer,
                                timestamp: new Date().toISOString()
                            });
                        });
                    }
                };
                
                // Expose the switch function globally
                window.switchToMainTab = function() {
                    // This will be overridden by the parent
                    console.log('Switching to main tab');
                };
            </script>
        </body>
        </html>
        """

    def _start_import_job(self, path: str, instr: str, do_summary: bool, do_flashcards: bool, do_quiz: bool):
        # show progress in the output area
        self.text_box.setHtml("<span style='color:#9aa0a6'>Importing and extracting text…</span>")
        QtWidgets.QApplication.processEvents()

        def worker():
            try:
                # 1) Extract text depending on file type
                if path.lower().endswith(".pptx"):
                    deck_text = self._extract_text_from_pptx(path)
                elif path.lower().endswith(".pdf"):
                    deck_text = self._extract_text_from_pdf(path)
                else:
                    deck_text = ""

                if not deck_text:
                    QtCore.QMetaObject.invokeMethod(
                        self.text_box, "setHtml",
                        QtCore.Qt.ConnectionType.QueuedConnection,
                        QtCore.Q_ARG(str, "<span style='color:#fca5a5'>No extractable text found.</span>")
                    )
                    return

                synopsis = ""
                # 2) Summary (optional)
                if do_summary:
                    # build the summary prompt, honoring user instruction if provided
                    user_note = f"\n\n# USER INSTRUCTION (highest priority):\n{instr}\n" if instr else ""
                    synopsis_prompt = (
                        "Create a concise, exam-oriented synthesis for a medical student.\n"
                        "• 12–18 short bullets total, grouped under 2–4 H2 headings (##).\n"
                        "• Prefer named tests/drugs, first-line items, classic signs, and key cutoffs.\n"
                        "• Avoid fluff; no slide titles."
                        f"{user_note}\n\n"
                        f"{deck_text}"
                    )
                    try:
                        resp = openai.ChatCompletion.create(
                            model="gpt-4o-mini",
                            messages=[
                                {"role":"system","content":"You produce tight, exam-ready syntheses."},
                                {"role":"user","content":synopsis_prompt}
                            ],
                            temperature=0.3
                        )
                        synopsis = (resp["choices"][0]["message"]["content"] or "").strip()
                    except Exception as e:
                        synopsis = "Import succeeded, but synthesis failed."

                    QtCore.QMetaObject.invokeMethod(
                        self.text_box, "setHtml",
                        QtCore.Qt.ConnectionType.QueuedConnection,
                        QtCore.Q_ARG(str, markdown_to_html(synopsis))
                    )
                    self.latest_context = synopsis

                # 3) Flashcards (optional) — reuse your function, but pass instruction
                if do_flashcards:
                    try:
                        self.generate_flashcards_from_gpt(deck_text, extra_instruction=instr or None)
                    except Exception as e:
                        print("Flashcard generation from slides failed:", e)

                # 4) Quiz (optional) — MUST be opened on UI thread
                if do_quiz:
                    QtCore.QTimer.singleShot(
                        0,
                        lambda bt=(synopsis or deck_text): self._open_quiz_dialog_from_text(bt)
                    )

            except Exception as e:
                QtCore.QMetaObject.invokeMethod(
                    self.text_box, "setHtml",
                    QtCore.Qt.ConnectionType.QueuedConnection,
                    QtCore.Q_ARG(str, f"<span style='color:#fca5a5'>Import failed: {e}</span>")
                )

        threading.Thread(target=worker, daemon=True).start()

        
    def _extract_text_from_pptx(self, path: str) -> str:
        if Presentation is None:
            raise RuntimeError("python-pptx not installed. Install with: pip install python-pptx")
        prs = Presentation(path)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            lines = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    for p in getattr(shape.text_frame, "paragraphs", []):
                        line = "".join(run.text for run in getattr(p, "runs", [])).strip()
                        if line:
                            lines.append(line)
            if lines:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(lines))
        return "\n\n".join(slides_text).strip()

    def _extract_text_from_pdf(self, path: str) -> str:
        if fitz is None:
            raise RuntimeError("PyMuPDF not installed. Install with: pip install pymupdf")
        doc = fitz.open(path)
        pages = []
        for i, page in enumerate(doc, 1):
            txt = page.get_text("text").strip()
            if txt:
                pages.append(f"[Page {i}]\n{txt}")
        doc.close()
        return "\n\n".join(pages).strip()
    
    def import_slides(self):
        """
        Pick a .pptx or .pdf, extract text, then:
        • Show a concise synthesis in text_box
        • Generate flashcards from the whole deck
        • Offer a quick 5-question MCQ quiz
        """
        path, _ = QtWidgets.QFileDialog.getOpenFileName(
            self, "Import Slides", "", "PowerPoint (*.pptx);;PDF (*.pdf)"
        )
        if not path:
            return

        self.text_box.setHtml("<span style='color:#9aa0a6'>Importing and extracting text…</span>")
        QtWidgets.QApplication.processEvents()

        def worker():
            try:
                if path.lower().endswith(".pptx"):
                    deck_text = self._extract_text_from_pptx(path)
                elif path.lower().endswith(".pdf"):
                    deck_text = self._extract_text_from_pdf(path)
                else:
                    deck_text = ""

                if not deck_text:
                    QtCore.QMetaObject.invokeMethod(
                        self.text_box, "setHtml", QtCore.Qt.ConnectionType.QueuedConnection,
                        QtCore.Q_ARG(str, "<span style='color:#fca5a5'>No extractable text found.</span>")
                    )
                    return

                # Synthesis
                synopsis_prompt = (
                    "Create a concise, exam-oriented synthesis of the following deck text for a medical student.\n"
                    "• 12–18 short bullets total, grouped under 2–4 H2 headings (##).\n"
                    "• Prefer named tests/drugs, first-line items, classic signs, and key cutoffs.\n"
                    "• Avoid fluff; no slide titles. Only include important exam relevant information.\n\n"
                    f"{deck_text}"
                )
                try:
                    resp = openai.ChatCompletion.create(
                        model="gpt-4o-mini",
                        messages=[
                            {"role":"system","content":"You produce tight, exam-ready syntheses."},
                            {"role":"user","content":synopsis_prompt}
                        ],
                        temperature=0.3
                    )
                    synopsis = resp["choices"][0]["message"]["content"].strip()
                except Exception as e:
                    synopsis = "Import succeeded, but synthesis failed."

                # Render
                QtCore.QMetaObject.invokeMethod(
                    self.text_box, "setHtml", QtCore.Qt.ConnectionType.QueuedConnection,
                    QtCore.Q_ARG(str, markdown_to_html(synopsis))
                )
                self.latest_context = synopsis

                # Flashcards from the whole deck (reuses your existing generator)
                try:
                    self.generate_flashcards_from_gpt(deck_text)
                except Exception as e:
                    print("Flashcard generation from slides failed:", e)

                # Quick quiz (5 MCQs) from synopsis (fallback to deck_text)
               # 3) Offer a quick quiz from the UI thread
                QtCore.QTimer.singleShot(
                    0,
                    lambda bt=(synopsis or deck_text): self._open_quiz_dialog_from_text(bt)
                )


            except Exception as e:
                QtCore.QMetaObject.invokeMethod(
                    self.text_box, "setHtml", QtCore.Qt.ConnectionType.QueuedConnection,
                    QtCore.Q_ARG(str, f"<span style='color:#fca5a5'>Import failed: {e}</span>")
                )

        threading.Thread(target=worker, daemon=True).start()




    def open_settings(self):
        # Standalone mode - use local user data only
        # updated = get_user_data(self.user_data['email'])
        # if not updated:
        #     QtWidgets.QMessageBox.critical(self, "Error", "No user data found.")
        #     return
        # self.user_data = updated

        # Open the settings dialog
        dialog = SettingsDialog(self.user_data, self)
        dialog.down_arrow_checkbox.setChecked(self.keyboard_enabled)

        if dialog.exec():  # User clicked OK
            # Update the setting based on the checkbox
            self.keyboard_enabled = dialog.down_arrow_checkbox.isChecked()
            self.user_data["custom_prompt"] = dialog.prompt_input.toPlainText().strip()
            print(f"Custom Prompt Saved: {self.user_data['custom_prompt']}")
            try:
                with open("user.json", "w") as f:
                    json.dump(self.user_data, f)
            except Exception as e:
                print(f"Error saving user data: {e}")
            print(f"Down arrow shortcut is now {'enabled' if self.keyboard_enabled else 'disabled'}")
    
    def _open_quiz_dialog_from_text(self, base_text: str):
        """Runs entirely on the UI thread. Creates MCQs, opens the dialog, handles results."""
        try:
            mcqs = self._generate_mcqs(base_text, n=5)
            if not mcqs:
                return
            dlg = QuizDialog(mcqs, parent=self)
            result = dlg.exec()
            if result == QtWidgets.QDialog.DialogCode.Accepted:
                wrongs = dlg.get_wrongs()
                score  = dlg.get_score()
                # Add flashcards for wrong answers
                for w in wrongs:
                    q_text = w["question"]
                    correct = w["choices"][w["answer_index"]]
                    exp     = w.get("explanation","")
                    ans = f"{correct}" + (f" — {exp}" if exp else "")
                    self.flashcards.append((q_text, ans))
                    # safe UI update
                    QtCore.QMetaObject.invokeMethod(
                        self, "generate_list_item_safe",
                        QtCore.Qt.ConnectionType.QueuedConnection,
                        QtCore.Q_ARG(str, q_text),
                        QtCore.Q_ARG(str, ans)
                    )
                QtWidgets.QMessageBox.information(
                    self, "Quiz Complete",
                    f"Score: {score['correct']}/{score['total']}\n"
                    f"{len(wrongs)} flashcard(s) added for review."
                )
        except Exception as e:
            print("Quiz dialog error:", e)



    def enter_compact_mode(self):
        """
        Hide the full window and show the compact bar under the top notch.
        """
        try:
            self.compact_bar.update_from_main()
            self.compact_bar.position_top_center()
            self.compact_bar.show()
            self.hide()
        except Exception as e:
            print(f"[MiniDock] enter_compact_mode error: {e}")

    def exit_compact_mode(self):
        """
        Hide the compact bar and restore the full window.
        """
        try:
            self.compact_bar.hide()
            self.show()
            self.raise_()
            self.activateWindow()
        except Exception as e:
            print(f"[MiniDock] exit_compact_mode error: {e}")





    @QtCore.pyqtSlot()
    def manual_summarize(self):
        

        global user_data
        
        sync_premium_status(self.user_data)
        # Standalone mode - no request limits
        # if not self.user_data.get("is_premium") and not can_make_request(self.user_data):
        #     self.text_box.setPlainText("You have used your 15 free daily requests. Please upgrade your plan to continue.")
        #     return

        print("Manual summarize triggered.")
        self.text_box.clear()

        self.capture_and_summarize()

        # Standalone mode - no request counting
        # if not self.user_data.get("is_premium"):
        #     increment_request_count(self.user_data["email"], self.user_data)

    def clear_summary(self):
        self.text_box.clear()
        #self.label.setText("Summary cleared")

    def _should_include_context(self, q: str) -> bool:
        """
        Return True only if the question is short/underspecified or uses deictic terms
        that need slide/session context. Otherwise, return False to avoid context bleed.
        """
        ql = q.lower().strip()
        if not ql:
            return False
        # Deictic / underspecified triggers
        triggers = [
            "this", "that", "these", "those", "above", "below",
            "here", "there", "slide", "on this", "in this", "from this",
            "as above", "as shown", "context", "the above", "the image"
        ]
        if any(t in ql for t in triggers):
            return True
        # Very short or fragment-like questions typically need context
        if len(ql) < 40:
            return True
        return False


    
    def process_question(self, question):
        try:
            q = (question or "").strip()
            if not q:
                return

            # Strict, concise answering rules to avoid rambling and stale context
            system_msg = (
                "You are a clinical tutor. Be precise and concise.\n"
                "Hard rules:\n"
                "• If the query expects a list (treatment/causes/investigations/risks/complications, etc.), return 3–6 SPECIFIC named items, comma-separated. No fluff.\n"
                "• Otherwise answer in 1–3 short sentences max, with concrete specifics only when essential.\n"
                "• Do NOT mention or rely on any prior slides or earlier topics unless explicit CONTEXT is provided to you below.\n"
                "• If the question is under-specified, answer the most common interpretation clinically; if ambiguity remains, ask ONE short clarifying question.\n"
                "• No preamble, no headings, no meta-commentary."
            )

            messages = [{"role": "system", "content": system_msg}]

            # Context: include ONLY when the question likely refers to 'this/that/above/slide' or is too short.
            try:
                need_ctx = self._should_include_context(q)
            except Exception:
                need_ctx = False

            if need_ctx:
                ctx = (getattr(self, "latest_context", "") or "").strip()
                if ctx:
                    messages.append({
                        "role": "system",
                        "content": "CONTEXT (use only if it disambiguates the user's short/underspecified query):\n" + ctx[:6000]
                    })

            # (Optional) lightweight recent history to keep a short thread without causing bleed
            # Keep at most the last 2 Q/A pairs.
            try:
                tail = self.conversation_history[-4:] if len(self.conversation_history) > 4 else self.conversation_history
                for turn in tail:
                    if isinstance(turn, dict) and "role" in turn and "content" in turn:
                        messages.append(turn)
            except Exception:
                pass

            # Current user question
            messages.append({"role": "user", "content": q})

            response = openai.ChatCompletion.create(
                model="gpt-4o-mini",
                messages=messages,
                temperature=0.2,      # tighter, more consistent
                max_tokens=350,       # hard cap to keep it tight
                stream=True
            )

            answer = ""
            for chunk in response:
                if 'choices' in chunk and 'delta' in chunk['choices'][0]:
                    piece = chunk['choices'][0]['delta'].get('content', '')
                    answer += piece

            answer = answer.strip()

            # Render
            QtCore.QMetaObject.invokeMethod(
                self.text_box,
                "setHtml",
                QtCore.Qt.ConnectionType.QueuedConnection,
                QtCore.Q_ARG(str, markdown_to_html(answer))
            )

            # Maintain compact history (at most last 4 turns total)
            try:
                self.conversation_history.append({"role": "user", "content": q})
                self.conversation_history.append({"role": "assistant", "content": answer})
                if len(self.conversation_history) > 8:
                    self.conversation_history = self.conversation_history[-8:]
            except Exception:
                pass

        except requests.exceptions.RequestException:
            self.text_box.setPlainText("❌ No internet connection. Please check your network and try again.")
        except Exception as e:
            self.text_box.setPlainText(f"Error: {e}")


    def _inline_ask(self):
        q = self.ask_input.text().strip()
        if not q:
            return

        # Immediately clear and show Processing…
        self.ask_input.clear()
        self.ask_input.setReadOnly(True)
        self.ask_input.setPlaceholderText("Processing…")
        QtWidgets.QApplication.processEvents()

        # Run the answer pipeline (synchronous in current code)
        try:
            self.process_question(q)
        finally:
            # Restore input box to idle state
            self.ask_input.setReadOnly(False)
            self.ask_input.clear()
            self.ask_input.setPlaceholderText("Ask anything…")
            # keep focus for fast follow-ups
            self.ask_input.setFocus(QtCore.Qt.FocusReason.ActiveWindowFocusReason)

    def open_screen_recording_settings(self):
        """Open macOS System Settings to Screen Recording permissions page."""
        if sys.platform == "darwin":
            try:
                # Try URL scheme for new System Settings (macOS Ventura 13+)
                subprocess.run([
                    "open", 
                    "x-apple.systempreferences:com.apple.preference.security?Privacy_ScreenCapture"
                ], check=False, timeout=2)
                print("📱 Opened System Settings → Privacy & Security → Screen Recording")
            except Exception:
                try:
                    # Fallback: Use AppleScript for older macOS or if URL scheme fails
                    applescript = '''
                    tell application "System Settings"
                        activate
                        reveal anchor "Privacy_ScreenCapture" of pane id "com.apple.preference.security"
                    end tell
                    '''
                    subprocess.run(["osascript", "-e", applescript], check=False, timeout=2)
                    print("📱 Opened System Settings via AppleScript")
                except Exception:
                    # Final fallback: just open System Settings
                    try:
                        subprocess.run(["open", "-b", "com.apple.systempreferences"], check=False, timeout=2)
                        print("📱 Opened System Settings (manual navigation required)")
                    except Exception as e:
                        print(f"Could not open settings: {e}")

    def open_accessibility_settings(self):
        """Open macOS System Settings to Accessibility permissions page."""
        if sys.platform == "darwin":
            try:
                # Try URL scheme for new System Settings (macOS Ventura 13+)
                subprocess.run([
                    "open", 
                    "x-apple.systempreferences:com.apple.preference.security?Privacy_Accessibility"
                ], check=False, timeout=2)
                print("📱 Opened System Settings → Privacy & Security → Accessibility")
            except Exception:
                try:
                    # Fallback: Use AppleScript
                    applescript = '''
                    tell application "System Settings"
                        activate
                        reveal anchor "Privacy_Accessibility" of pane id "com.apple.preference.security"
                    end tell
                    '''
                    subprocess.run(["osascript", "-e", applescript], check=False, timeout=2)
                    print("📱 Opened Accessibility Settings via AppleScript")
                except Exception:
                    # Final fallback: just open System Settings
                    try:
                        subprocess.run(["open", "-b", "com.apple.systempreferences"], check=False, timeout=2)
                        print("📱 Opened System Settings (manual navigation required)")
                    except Exception as e:
                        print(f"Could not open settings: {e}")

    def check_accessibility_permission(self):
        """Check if Accessibility permission is granted (macOS)."""
        if sys.platform != "darwin":
            return True
        try:
            # Use AppleScript to check if we have accessibility permissions
            script = '''
            tell application "System Events"
                return UI elements enabled
            end tell
            '''
            result = subprocess.run(
                ["osascript", "-e", script],
                capture_output=True,
                text=True,
                timeout=2
            )
            return "true" in result.stdout.lower()
        except Exception:
            # If we can't check, assume we need to request
            return False

    def check_screen_permissions(self):
        """Test if screen recording and accessibility permissions are granted (macOS)."""
        if sys.platform != "darwin":
            return
        
        screen_ok = False
        accessibility_ok = self.check_accessibility_permission()
        
        # Check screen recording
        try:
            with mss.mss() as sct:
                sct.grab(sct.monitors[0])
            screen_ok = True
            print("✅ Screen recording permissions OK")
        except Exception as e:
            error_msg = str(e).lower()
            if "permission" in error_msg or "access" in error_msg:
                screen_ok = False
        
        # Check what's missing
        missing = []
        if not screen_ok:
            missing.append("Screen Recording")
        if not accessibility_ok:
            missing.append("Accessibility")
        
        if missing:
            missing_str = " and ".join(missing)
            permission_html = f"""
            <div style='padding: 12px; color: #ff9800; font-size: 13px; line-height: 1.5; background-color: #2a1e00; border-radius: 8px; border: 1px solid #ff9800;'>
                <strong style='font-size: 14px;'>⚠️ macOS Permissions Required</strong>
                <p style='margin: 8px 0 4px 0;'>Akson needs <strong>{missing_str}</strong> permission(s) to capture and process your screen.</p>
                <p style='margin: 8px 0 4px 0;'><strong>Steps:</strong></p>
                <ol style='margin: 4px 0 8px 20px; padding-left: 4px;'>
                    <li>System Settings will open automatically</li>
                    <li>Find <strong>Python</strong> (or Terminal/Python3) in the permission list</li>
                    <li>Enable the toggle(s) next to it</li>
                    <li>Restart Akson after enabling</li>
                </ol>
                <p style='margin: 8px 0 0 0; color: #ccc; font-size: 12px;'>
                    <em>💡 Go to: System Settings → Privacy & Security → {' → '.join(missing)}</em>
                </p>
            </div>
            """
            self.text_box.setHtml(permission_html)
            print(f"⚠️ Missing permissions: {missing_str}")
            # Auto-open settings - open both if both are missing
            QtCore.QTimer.singleShot(500, self.open_screen_recording_settings if not screen_ok else self.open_accessibility_settings)
            if not screen_ok and not accessibility_ok:
                # Open both with a delay between them
                QtCore.QTimer.singleShot(1000, self.open_accessibility_settings)
    
    def _setup_app_icon(self):
        """Set app icon for macOS dock and window."""
        if sys.platform == "darwin":
            icon_path = os.path.join(os.path.dirname(__file__), "icons", "akson.png")
            if os.path.exists(icon_path):
                self.setWindowIcon(QIcon(icon_path))
                # Also set app icon for dock
                app = QApplication.instance()
                if app:
                    app.setWindowIcon(QIcon(icon_path))
    
    def _create_menu_bar(self):
        """Create macOS-style menu bar."""
        if sys.platform != "darwin":
            return
        
        menubar = self.menuBar()
        
        # File Menu
        file_menu = menubar.addMenu("File")
        new_action = QAction("New", self)
        new_action.setShortcut(QKeySequence("Ctrl+N"))
        file_menu.addAction(new_action)
        
        open_action = QAction("Open...", self)
        open_action.setShortcut(QKeySequence("Ctrl+O"))
        file_menu.addAction(open_action)
        
        file_menu.addSeparator()
        
        close_action = QAction("Close Window", self)
        close_action.setShortcut(QKeySequence("Ctrl+W"))
        close_action.triggered.connect(self.close)
        file_menu.addAction(close_action)
        
        file_menu.addSeparator()
        
        quit_action = QAction("Quit Akson", self)
        quit_action.setShortcut(QKeySequence("Ctrl+Q"))
        quit_action.triggered.connect(QApplication.instance().quit)
        file_menu.addAction(quit_action)
        
        # Edit Menu
        edit_menu = menubar.addMenu("Edit")
        
        undo_action = QAction("Undo", self)
        undo_action.setShortcut(QKeySequence("Ctrl+Z"))
        edit_menu.addAction(undo_action)
        
        redo_action = QAction("Redo", self)
        redo_action.setShortcut(QKeySequence("Ctrl+Shift+Z"))
        edit_menu.addAction(redo_action)
        
        edit_menu.addSeparator()
        
        cut_action = QAction("Cut", self)
        cut_action.setShortcut(QKeySequence("Ctrl+X"))
        cut_action.triggered.connect(lambda: self._handle_clipboard_action("cut"))
        edit_menu.addAction(cut_action)
        
        copy_action = QAction("Copy", self)
        copy_action.setShortcut(QKeySequence("Ctrl+C"))
        copy_action.triggered.connect(lambda: self._handle_clipboard_action("copy"))
        edit_menu.addAction(copy_action)
        
        paste_action = QAction("Paste", self)
        paste_action.setShortcut(QKeySequence("Ctrl+V"))
        paste_action.triggered.connect(lambda: self._handle_clipboard_action("paste"))
        edit_menu.addAction(paste_action)
        
        edit_menu.addSeparator()
        
        select_all_action = QAction("Select All", self)
        select_all_action.setShortcut(QKeySequence("Ctrl+A"))
        select_all_action.triggered.connect(lambda: self._handle_clipboard_action("select_all"))
        edit_menu.addAction(select_all_action)
        
        # View Menu
        view_menu = menubar.addMenu("View")
        
        # Window Menu
        window_menu = menubar.addMenu("Window")
        
        minimize_action = QAction("Minimize", self)
        minimize_action.setShortcut(QKeySequence("Ctrl+M"))
        minimize_action.triggered.connect(self.showMinimized)
        window_menu.addAction(minimize_action)
        
        zoom_action = QAction("Zoom", self)
        zoom_action.triggered.connect(self._toggle_maximize)
        window_menu.addAction(zoom_action)
        
        window_menu.addSeparator()
        
        bring_all_action = QAction("Bring All to Front", self)
        bring_all_action.triggered.connect(self.raise_)
        window_menu.addAction(bring_all_action)
        
        # Help Menu
        help_menu = menubar.addMenu("Help")
        
        about_action = QAction("About Akson", self)
        about_action.triggered.connect(self._show_about)
        help_menu.addAction(about_action)
        
        shortcuts_action = QAction("Keyboard Shortcuts", self)
        shortcuts_action.setShortcut(QKeySequence("Ctrl+?"))
        shortcuts_action.triggered.connect(self._show_shortcuts)
        help_menu.addAction(shortcuts_action)
    
    def _setup_keyboard_shortcuts(self):
        """Setup standard macOS keyboard shortcuts."""
        # Cmd+H (Hide)
        hide_action = QAction(self)
        hide_action.setShortcut(QKeySequence("Ctrl+H"))
        hide_action.triggered.connect(self.hide)
        self.addAction(hide_action)
        
        # Cmd+M (Minimize) - already in menu but add here too
        minimize_action = QAction(self)
        minimize_action.setShortcut(QKeySequence("Ctrl+M"))
        minimize_action.triggered.connect(self.showMinimized)
        self.addAction(minimize_action)
    
    def _setup_macos_theme_detection(self):
        """Detect macOS system theme and sync with webview."""
        if sys.platform != "darwin":
            return
        
        try:
            # Use PyQt6's native theme detection
            app = QApplication.instance()
            if app:
                # Get system palette to detect theme
                palette = app.palette()
                bg_color = palette.color(palette.ColorRole.Window)
                # Dark mode typically has darker background
                is_dark = bg_color.lightness() < 128
                system_theme = 'dark' if is_dark else 'light'
                
                # Store for later use
                self._system_theme = system_theme
                
                # Send to webview after a delay to ensure it's ready
                QtCore.QTimer.singleShot(500, lambda: self._send_theme_to_webview(system_theme))
                
                print(f"🌓 Detected macOS system theme: {system_theme}")
        except Exception as e:
            print(f"⚠️  Error detecting system theme: {e}")
            self._system_theme = 'dark'  # Default
    
    def _send_theme_to_webview(self, theme):
        """Send theme to webview (slides_working.py) to sync sidebar and PDF viewer."""
        try:
            # This will be handled by the webview's message handler
            # The webview should listen for postMessage with theme info
            pass  # Implementation depends on how webview communication works
        except Exception as e:
            print(f"⚠️  Error sending theme to webview: {e}")
    
    def _handle_clipboard_action(self, action):
        """Handle clipboard actions from menu."""
        clipboard = QApplication.clipboard()
        if action == "copy":
            # Get selected text from focused widget
            widget = QApplication.focusWidget()
            if hasattr(widget, 'selectedText'):
                text = widget.selectedText()
                if text:
                    clipboard.setText(text)
        elif action == "paste":
            widget = QApplication.focusWidget()
            if hasattr(widget, 'paste'):
                widget.paste()
        elif action == "cut":
            widget = QApplication.focusWidget()
            if hasattr(widget, 'cut'):
                widget.cut()
        elif action == "select_all":
            widget = QApplication.focusWidget()
            if hasattr(widget, 'selectAll'):
                widget.selectAll()
    
    def _toggle_maximize(self):
        """Toggle window maximize/restore."""
        if self.isMaximized():
            self.showNormal()
        else:
            self.showMaximized()
    
    def _show_about(self):
        """Show About dialog."""
        QtWidgets.QMessageBox.about(
            self,
            "About Akson",
            "Akson - Study Co-pilot\n\n"
            "An AI-powered study assistant for medical students.\n\n"
            "Version 1.0"
        )
    
    def _show_shortcuts(self):
        """Show keyboard shortcuts dialog."""
        shortcuts_text = """
Keyboard Shortcuts:

File:
  Cmd+N          New
  Cmd+O          Open
  Cmd+W          Close Window
  Cmd+Q          Quit

Edit:
  Cmd+Z          Undo
  Cmd+Shift+Z    Redo
  Cmd+X          Cut
  Cmd+C          Copy
  Cmd+V          Paste
  Cmd+A          Select All

Window:
  Cmd+M          Minimize
  Cmd+H          Hide

Help:
  Cmd+?          Keyboard Shortcuts
        """
        QtWidgets.QMessageBox.information(self, "Keyboard Shortcuts", shortcuts_text.strip())

    def capture_and_summarize(self):
        print(f"\n{'='*60}")
        print("🚀 AKSON: Starting screen capture and OCR process")
        print(f"{'='*60}\n")
        
        # Quick permission check on macOS
        if sys.platform == "darwin":
            screen_ok = False
            try:
                with mss.mss() as sct:
                    sct.grab(sct.monitors[0])
                screen_ok = True
                print("✅ Screen Recording permission: OK")
            except Exception as e:
                print(f"⚠️  Screen Recording permission: FAILED ({str(e)[:50]})")
            
            accessibility_ok = self.check_accessibility_permission()
            if accessibility_ok:
                print("✅ Accessibility permission: OK")
            else:
                print("⚠️  Accessibility permission: FAILED")
        
        self.loading_label.setText("Processing selection…")
        self.loading_label.show()
        try:
            # 1) Clear any old summary so it's not in the shot
            self.text_box.clear()
            # Do NOT hide the window; capture proceeds while visible per user preference
            QtWidgets.QApplication.processEvents()
            
            try:
                with mss.mss() as sct:
                    print("📷 Capturing screen...")
                    screenshot = sct.grab(sct.monitors[0])
                    img = np.array(screenshot)
                    img = cv2.cvtColor(img, cv2.COLOR_BGRA2BGR)
                    print("✅ Screen captured successfully!")
            except Exception as screen_error:
                error_msg = str(screen_error).lower()
                # Restore window before handling errors
                try:
                    self.show()
                    self.raise_()
                    self.activateWindow()
                except Exception:
                    pass
                self.show()
                if sys.platform == "darwin" and ("permission" in error_msg or "access" in error_msg or "screen recording" in error_msg):
                    # Check both permissions
                    accessibility_ok = self.check_accessibility_permission()
                    missing = ["Screen Recording"]
                    if not accessibility_ok:
                        missing.append("Accessibility")
                    missing_str = " and ".join(missing)
                    
                    permission_html = f"""
                    <div style='padding: 12px; color: #ff9800; font-size: 13px; line-height: 1.5; background-color: #2a1e00; border-radius: 8px; border: 1px solid #ff9800;'>
                        <strong style='font-size: 14px;'>⚠️ macOS Permissions Required</strong>
                        <p style='margin: 8px 0 4px 0;'>Akson needs <strong>{missing_str}</strong> permission(s) to work properly.</p>
                        <p style='margin: 8px 0 4px 0;'><strong>Steps:</strong></p>
                        <ol style='margin: 4px 0 8px 20px; padding-left: 4px;'>
                            <li>System Settings will open automatically</li>
                            <li>Find <strong>Python</strong> (or Terminal/Python3) in the permission list(s)</li>
                            <li>Enable the toggle(s) next to it</li>
                            <li>Restart Akson after enabling</li>
                        </ol>
                        <p style='margin: 8px 0 0 0; color: #ccc; font-size: 12px;'>
                            <em>💡 Go to: System Settings → Privacy & Security → {' → '.join(missing)}</em>
                        </p>
                    </div>
                    """
                    self.text_box.setHtml(permission_html)
                    self.loading_label.hide()
                    print(f"❌ Permission error: {missing_str}")
                    # Auto-open settings - open both if needed
                    QtCore.QTimer.singleShot(500, self.open_screen_recording_settings)
                    if not accessibility_ok:
                        QtCore.QTimer.singleShot(1000, self.open_accessibility_settings)
                    return
                else:
                    # Other error - show full details
                    import html, traceback
                    error_details = traceback.format_exc()
                    print(f"\n{'='*60}")
                    print("❌ SCREEN CAPTURE ERROR")
                    print(f"{'='*60}")
                    print(f"Error Type: {type(screen_error).__name__}")
                    print(f"Error Message: {screen_error}")
                    print(f"\nFull Traceback:")
                    print(error_details)
                    print(f"{'='*60}\n")
                    
                    error_html = f"""
                    <div style='padding: 12px; color: #f44336; font-size: 13px; line-height: 1.5; background-color: #2a0000; border-radius: 8px; border: 1px solid #f44336;'>
                        <strong style='font-size: 14px;'>❌ Screen Capture Error</strong>
                        <p style='margin: 8px 0 4px 0;'><strong>Error:</strong> {html.escape(str(screen_error))}</p>
                        <p style='margin: 4px 0; color: #ccc; font-size: 12px;'>
                            Check the console for full error details.
                        </p>
                    </div>
                    """
                    self.text_box.setHtml(error_html)
                    self.loading_label.hide()
                    return

            print(f"\n{'='*60}")
            print(f"📸 SCREENSHOT CAPTURED: {img.shape[1]}x{img.shape[0]} pixels")
            print(f"{'='*60}")
            cv2.imwrite("screenshot.png", img)
            print("💾 Screenshot saved to: screenshot.png")
            # Restore the window after capture
            try:
                self.show()
                self.raise_()
                self.activateWindow()
            except Exception:
                pass
            
            print("\n🔍 Running OCR with Tesseract...")
            extracted_text = ""
            try:
                # Check if Tesseract is available
                try:
                    pytesseract.get_tesseract_version()
                    print("✅ Tesseract found")
                except Exception as tesseract_check_error:
                    print(f"⚠️  Tesseract check warning: {tesseract_check_error}")
                
                extracted_text = pytesseract.image_to_string(Image.fromarray(img))
                text_length = len(extracted_text.strip())
                
                print(f"\n{'='*60}")
                print("✅ OCR EXTRACTION COMPLETE")
                print(f"{'='*60}")
                print(f"📊 Extracted text length: {text_length} characters")
                print(f"📄 Number of lines: {len(extracted_text.strip().splitlines())}")
                print(f"\n{'─'*60}")
                print("📝 EXTRACTED TEXT:")
                print(f"{'─'*60}")
                if extracted_text.strip():
                    print(extracted_text.strip())
                else:
                    print("⚠️  NO TEXT DETECTED - Screenshot may be empty or text unclear")
                print(f"{'─'*60}\n")
                
                self.last_extracted_text = extracted_text.strip()
            except Exception as ocr_error:
                import traceback
                error_details = traceback.format_exc()
                print(f"\n{'='*60}")
                print("❌ OCR ERROR")
                print(f"{'='*60}")
                print(f"Error Type: {type(ocr_error).__name__}")
                print(f"Error Message: {ocr_error}")
                print(f"\nFull Traceback:")
                print(error_details)
                print(f"{'='*60}\n")
                
                import html
                ocr_error_html = f"""
                <div style='padding: 12px; color: #f44336; font-size: 13px; line-height: 1.5; background-color: #2a0000; border-radius: 8px; border: 1px solid #f44336;'>
                    <strong style='font-size: 14px;'>❌ OCR Error</strong>
                    <p style='margin: 8px 0 4px 0;'><strong>Error:</strong> {html.escape(str(ocr_error))}</p>
                    <p style='margin: 4px 0; color: #ccc; font-size: 12px;'>
                        Check the console for full error details. This might indicate:<br/>
                        • Tesseract OCR is not installed<br/>
                        • Tesseract path is not configured<br/>
                        • Image format issue
                    </p>
                </div>
                """
                self.text_box.setHtml(ocr_error_html)
                self.loading_label.hide()
                self.last_extracted_text = ""
                return
            self.loading_label.setText("Processing selection…")
            self.loading_label.show()
            if extracted_text.strip():
                selected_mode = self.mode_selector.currentText()

                if selected_mode == "Explain":
                    prompt = f"""You are Akson, a medical study co-pilot and clinical educator. Your job is to convert a raw medical lecture slide into a high-yield, deeply explained, clinically relevant, exam-ready teaching resource. You're not summarizing — you're teaching.

🧠 CORE INSTRUCTIONS
If no text is selected (by selected i mean the user has clearly highlighted certain text):
1️⃣ Read & Filter Smartly
Carefully read the entire slide. Completely ignore accidental, irrelevant, or non-medical text. Focus only on medically important content.

2️⃣ Teach, Don't Rephrase
Do not paraphrase or simplify blindly. Instead, explain everything clearly and intelligently — as if teaching a bright but uncertain junior medical student. For every concept, mechanism, pathway, enzyme, receptor, disease, drug, abbreviation (e.g. HMG-CoA, ACC), or instruction:

Define it (briefly and clearly)
Explain what it does
Explain why it matters clinically
Add what the slide leaves out but a good teacher would include
Assume the student has knowledge gaps. Fill them proactively.

3️⃣ Fill In What the Slide Leaves Out
If the slide vaguely says things like "laboratory tests," "diagnosis," or "treatment," you must:

List what those tests might be
Briefly explain why they're ordered
Include standard options (e.g. FBC, ECG, CRP, CT) where appropriate
Never allow vague or incomplete statements to pass through unexplained.
4️⃣ Always Add Clinical Context
Every explanation must tie into real-world patient care. Show how the content connects to:

Symptoms and red flags
Clinical reasoning and diagnosis
Investigations and imaging
Treatments (first-line vs. second-line)
Prognosis
Disease mechanisms and complications
5️⃣ Prioritise High-Yield Content
Mark 🔑 high-yield facts, 💡 mnemonics or tips, and ⚠️ exam traps or common misconceptions.
Mention anything that's commonly tested or clinically critical. Highlight contrasts and "commonly confused" points.

6️⃣ Use Active Recall Sparingly
After each major idea or section, insert 1 short active recall question (if appropriate).
Example: 💭 "What's the main risk factor for this condition?"
Keep questions simple and relevant.

7️⃣ Structure for Maximum Learning
Use a clean, consistent layout:

Clear section headings
Bullet points
Line spacing for readability
Concise but precise medical language
Avoid waffle. Use emojis (🧠, 💉, ⚠️, ❤️) sparingly to boost engagement.
8️⃣ Always Add Value Beyond the Slide
If the slide is shallow or incomplete, you must enrich it. Fill in mechanisms. Suggest simple mnemonics. Clarify unclear pathways. Break complex concepts into digestible steps. Think: "If the student sees this in a viva or ward round, what would they need to say or understand?"

9️⃣ Your Mental Model
You are a sharp, clinical, exam-aware tutor.
You are here to build true understanding, not copy text.
The final output must help the student:
x
Grasp the concept
Retain it
Apply it
Pass their exams with confidence

If TEXT IS SELECTED or hovered over or highlighted:
only explain the word selected, or the phrase or sentence(s) selected.


        {extracted_text.strip()}

        """
                else:
                  
                    prompt = f"""You summarise OCR text from medical lecture slides into SHORT, high-yield bullets for exam revision.

GOAL
Return a compact list of independent facts in arrow/abbreviation style (memory notes), not prose.

HARD RULES (must obey)
- Use emojis to boost engagement and help with navigation. around 2-3 per slide.
 - start with a heading so we know what you will speak about.
  - Make it short but easy to understand. use words and wording that makes it easy to comprehend
  - Use bold text where needed, use italics and emojis also just so that the overall output is easy to navigate and read at a glance. 
. 
- 4–6 bullets total. One fact per bullet. No paragraphs.
- Do NOT start bullets with category labels (e.g., "Epidemiology:", "Prevalence:", "Demographics:", "Skeletal destruction:").
  • Allowed labels only when logically required: **Criteria** and **Comparison**.
- Never chain multiple facts with ";" or long commas. Split into separate bullets.
- Use symbols & shorthand aggressively:
  ↑ increased/elevated, ↓ decreased/low, → leads to/causes, ~ approx, ≈ ratio, κ/λ light chains.
- Prefer standard medical abbreviations:
  myocardial infarction → **MI**
  myocardial ischemia → **ischemia**
  ST elevation → **ST ↑**, ST depression → **ST ↓**, T-wave inversion(s) → **TWI**
  ST-elevation MI → **STEMI**, non-ST-elevation MI → **NSTEMI**
  unstable angina → **UA**, stable angina → **SA**
  electrocardiogram → **ECG**
- Ratios compact: "M:F ≈ 3:1", "Blacks:Whites ≈ 2:1".
- Minimal bold: only the key term in a bullet (e.g., **NSTEMI**, **Troponin I**, **Russell bodies**).
- Zero filler or basic definitions.
- Not every bullet point needs a subheading! You do not need to say narrowing -> progressive narrowing of aortic valve orifive, for example, instead just say - progressive narrowing of... thats all the bullet point needs to be THERE DOESNT HAVE TO BE A HEADING FOR EACH BULLET POINT

REASONING (internal; do not output)
1) Identify intent (definition/essence, criteria, comparison, atypical features, pathology/morphology, epidemiology).
Identify also what the slide is talking about. just because a word is mentioned doesnt mean the slide is about that topic. see what the slide is talking about and match it to what logically should be the most appropriate title and use that as your heading. For example, if 'Heart failure with reduced ejection fraction is mentioned but acute coronary syndrome is also mentioned and you are confused what the title is, you should see what the majority of the text lines up with.!!! this is important to get right.
2) Select the 4–6 highest-yield, exam-testable facts; prefer differentiators and criteria.
3) Rewrite each as: **Label** → compressed fact(s) with symbols/abbrevs.
4) Self-check before output:
   - Bullet count 4–6 or as little as possible.
   - No bullet starts with a category label (except **Criteria**/**Comparison** when necessary).
   - No ";" inside bullets.
   - Each bullet ≤ ~18–22 words.

FORMATS WHEN NEEDED
- **Criteria:** Use a single bullet if it's the classic triad/"≥2 of 3". Otherwise split into bullets.
- **Comparison:** Up to 3 bullets, one per entity, each starting with a bold entity then its defining line.

Return ONLY the bullets.


{extracted_text.strip()}

"""
                

                
                custom_prompt = self.user_data.get("custom_prompt", "")
                system_msg = "You are a helpful medical tutor."
                if custom_prompt:
                    system_msg += (
                        "\n\n# USER OVERRIDE INSTRUCTIONS (highest priority):\n"
                        f"{custom_prompt}\n"
                        "If any instruction above conflicts with previous instructions, "
                        "FOLLOW THE USER OVERRIDE INSTRUCTIONS."
                    )

                response = openai.ChatCompletion.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": system_msg},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.4,
                    stream=True
                )


                summary = ""
                

                for chunk in response:
                    if 'choices' in chunk and 'delta' in chunk['choices'][0]:
                        content_piece = chunk['choices'][0]['delta'].get('content', '')
                        summary += content_piece
                self.latest_context = summary.strip()

                QtCore.QMetaObject.invokeMethod(
                    self, "_collect_lecture_piece",
                    QtCore.Qt.ConnectionType.QueuedConnection,
                    QtCore.Q_ARG(str, summary.strip()),
                    QtCore.Q_ARG(str, selected_mode)
                )


                # After streaming fully finishes:
                
                QtCore.QTimer.singleShot(0, lambda: self.text_box.setHtml(markdown_to_html(summary.strip())))
                # Save this exchange into conversation memory
               

                #self.label.setText("Summary ready")


            else:
                #self.label.setText("No text detected.")
                self.text_box.setPlainText("No text detected")

        except Exception as e:
            print(f"Error in capture_and_summarize: {e}")
            #self.label.setText("Error.")
            self.text_box.setPlainText(f"Error: {e}")

    def _generate_mcqs(self, source_text: str, n: int = 5):
        """
        Returns a list of MCQs:
        [{"question": str, "choices": [A,B,C,D], "answer_index": int, "explanation": str}, ...]
        """
        sys = (
            "You write single-best-answer clinical MCQs for medical students. "
            "Use clinical vignettes and realistic distractors. Avoid trivia, avoid 'All of the above'. Make the questions like questions that would actually come up in medical school exams, completely based on the content provided."
            "Target diagnosis, best initial investigation, best next step, or first-line management. Exactly 4 options."
        )
        user = (
            f"Create {n} clinical MCQs from the content below.\n"
            "Constraints:\n"
            "- Short vignette (1–3 sentences) when useful.\n"
            "- 4 plausible options; one correct (answer_index 0–3).\n"
            "- Use named tests/drugs and sensible cutoffs if relevant.\n"
            "- Provide a one-sentence explanation that states *why* the correct option is best.\n"
            "Output STRICT JSON (list of objects) with keys: question, choices (array of 4 strings), answer_index, explanation.\n\n"
            f"CONTENT:\n{source_text}"
        )
        resp = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role":"system","content":sys},{"role":"user","content":user}],
            temperature=0.3
        )
        raw = resp["choices"][0]["message"]["content"]
        import json, re
        try:
            return json.loads(raw)
        except:
            m = re.search(r"```json\s*(.+?)\s*```", raw, re.S|re.I)
            if m:
                return json.loads(m.group(1))
        return []

    # === DRAGGABLE window ===
    def mousePressEvent(self, event):
        self.oldPos = event.globalPosition().toPoint()

    def mouseMoveEvent(self, event):
        delta = event.globalPosition().toPoint() - self.oldPos
        self.move(self.x() + delta.x(), self.y() + delta.y())
        self.oldPos = event.globalPosition().toPoint()
    
    def start_keyboard_listener(self):
        from pynput import keyboard

        self.key_press_times = deque()

        def on_press(key):
            if not self.keyboard_enabled:
                return  # shortcut disabled by user
            # rest of the code...
            try:
                # Filter: Only react to RIGHT and DOWN arrows
                if key not in {Key.right, Key.down}:
                    return  # Ignore all other keys

                now = time.time()
                self.key_press_times.append(now)

                # Remove key presses older than 5 seconds
                while self.key_press_times and now - self.key_press_times[0] > 5:
                    self.key_press_times.popleft()

                if len(self.key_press_times) > 2:
                    # User is spamming — reset timer
                    if hasattr(self, 'debounce_timer') and self.debounce_timer.isActive():
                        self.debounce_timer.stop()

                    self.debounce_timer = QtCore.QTimer()
                    self.debounce_timer.setSingleShot(True)
                    self.debounce_timer.timeout.connect(lambda: QtCore.QMetaObject.invokeMethod(
                        self, "manual_summarize", QtCore.Qt.ConnectionType.QueuedConnection))
                    self.debounce_timer.start(1500)  # wait 1.5s after last key press
                else:
                    # If not spamming, trigger immediately
                    QtCore.QMetaObject.invokeMethod(
                        self, "manual_summarize", QtCore.Qt.ConnectionType.QueuedConnection)
            except Exception as e:
                print(f"Keyboard listener error: {e}")

        listener = keyboard.Listener(
            on_press=on_press,
            suppress=False
        )
        listener.start()

    

    def start_lecture(self):
        filename, _ = QtWidgets.QFileDialog.getSaveFileName(self, "Save Anki Deck", "", "TSV Files (*.tsv)")
        if filename:
            if not filename.endswith(".tsv"):
                filename += ".tsv"
            self.anki_filename = filename
            self.flashcards = []
            self.lecture_pieces = []
            self.lecture_active = True

            self.create_flashcard_button.setEnabled(True)

            self.end_lecture_button.setEnabled(True)

            # tiny hint for the user
            self.text_box.setHtml("<span style='color:#9aa0a6'>Lecture started — summaries will be collected here while you work.</span>")

    def end_lecture(self):
        if not self.lecture_active:
            QtWidgets.QMessageBox.information(self, "End Lecture", "No active lecture.")
            return
        self.lecture_active = False
        self.end_lecture_button.setEnabled(False)

        if not self.lecture_pieces:
            QtWidgets.QMessageBox.information(self, "End Lecture", "No collected summaries.")
            return

        # 1) Build a concise master summary (AI) — same as before
        all_text = "\n\n---\n\n".join(p["text"] for p in self.lecture_pieces)
        try:
            prompt = (
                "You are Akson. Create a single concise, exam-ready synthesis of the entire session.\n"
                "• Group related ideas under clear H2 headings (## ).\n"
                "• 10–16 short bullets total (not per section). No fluff, no slide titles.\n"
                "• Prefer named items and cutoffs. Avoid repeating the same fact.\n\n"
                f"SESSION CONTENT:\n{all_text}"
            )
            resp = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role":"system","content":"You produce tight, exam-oriented summaries."},
                    {"role":"user","content":prompt}
                ],
                temperature=0.3
            )
            master_summary = resp["choices"][0]["message"]["content"].strip()
            self.latest_context = master_summary
            QtCore.QMetaObject.invokeMethod(
                self.text_box, "setHtml",
                QtCore.Qt.ConnectionType.QueuedConnection,
                QtCore.Q_ARG(str, markdown_to_html(master_summary))
            )
        except Exception as e:
            master_summary = ""
            QtWidgets.QMessageBox.warning(self, "Summary Error", f"Failed to create master summary:\n{e}")

        # 2) Generate ~5 MCQs and run the quiz BEFORE any export
        ran_quiz = False
        try:
            mcqs = self._generate_mcqs(master_summary or all_text, n=5)
            if mcqs:
                dlg = QuizDialog(mcqs, parent=self)
                result = dlg.exec()
                if result == QtWidgets.QDialog.DialogCode.Accepted:
                    wrongs = dlg.get_wrongs()
                    score  = dlg.get_score()
                    # Add flashcards for wrong answers
                    for w in wrongs:
                        q_text = w["question"]
                        correct = w["choices"][w["answer_index"]]
                        exp     = w.get("explanation","")
                        ans = f"{correct}" + (f" — {exp}" if exp else "")
                        self.flashcards.append((q_text, ans))
                        QtCore.QMetaObject.invokeMethod(
                            self, "generate_list_item_safe",
                            QtCore.Qt.ConnectionType.QueuedConnection,
                            QtCore.Q_ARG(str, q_text),
                            QtCore.Q_ARG(str, ans)
                        )
                    QtWidgets.QMessageBox.information(
                        self, "Quiz Complete",
                        f"Score: {score['correct']}/{score['total']}\n"
                        f"{len(wrongs)} flashcard(s) added for review."
                    )
                    ran_quiz = True
                else:
                    # If skipped, no cards added; proceed to exports
                    ran_quiz = False
        except Exception as e:
            QtWidgets.QMessageBox.warning(self, "Quiz Error", f"Failed to run quiz:\n{e}")

        # 3) Offer to export flashcards (TSV) NOW — after quiz so new cards are included
        if self.flashcards:
            do_export = QtWidgets.QMessageBox.question(
                self,
                "Export Deck",
                "Export your flashcards deck now (TSV)?",
                QtWidgets.QMessageBox.StandardButton.Yes | QtWidgets.QMessageBox.StandardButton.No,
                QtWidgets.QMessageBox.StandardButton.Yes
            )
            if do_export == QtWidgets.QMessageBox.StandardButton.Yes:
                self.export_deck()


        # 4) Export ALL individual summaries as a PDF (unchanged)
        try:
            pdf_path, _ = QtWidgets.QFileDialog.getSaveFileName(self, "Export Session Summaries (PDF)", "", "PDF (*.pdf)")
            if pdf_path:
                if not pdf_path.lower().endswith(".pdf"):
                    pdf_path += ".pdf"
                html_sections = []
                for i, p in enumerate(self.lecture_pieces, 1):
                    html_sections.append(
                        f"<h2 style='margin:12px 0 6px 0;'>Section {i} — {p['mode']}</h2>"
                        f"<div style='white-space:pre-wrap; line-height:1.55;'>{markdown_to_html(p['text'])}</div>"
                        "<hr style='border:none;border-top:1px solid #333;margin:12px 0;'>"
                    )
                full_html = (
                    "<html><head><meta charset='utf-8'>"
                    "<style>body{font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Roboto,sans-serif;color:#111;}"
                    "h1{font-size:20px;margin:0 0 12px 0;} h2{font-size:16px;margin:14px 0 8px 0;}"
                    "</style></head><body>"
                    "<h1>Akson — Lecture Summaries</h1>"
                    + "".join(html_sections) +
                    "</body></html>"
                )
                doc = _QtGui.QTextDocument()
                doc.setHtml(full_html)
                printer = QPrinter(QPrinter.PrinterMode.HighResolution)
                printer.setOutputFormat(QPrinter.OutputFormat.PdfFormat)
                printer.setOutputFileName(pdf_path)
                doc.print(printer)
        except Exception as e:
            QtWidgets.QMessageBox.warning(self, "PDF Error", f"Failed to export PDF:\n{e}")
        
        try:
            # Clear in-memory flashcards and UI list so they won't be exported/added twice
            self.flashcards.clear()
            self.flashcard_list.clear()

            # Reset target filename so a new lecture prompts for a fresh save path
            self.anki_filename = None

            # Disable flashcard creation until a new lecture starts
            self.create_flashcard_button.setEnabled(False)

        except Exception as e:
            QtWidgets.QMessageBox.warning(self, "Cleanup Warning", f"Post-lecture cleanup issue:\n{e}")

    @QtCore.pyqtSlot(str, str)
    def _collect_lecture_piece(self, text, mode):
        if self.lecture_active and text.strip():
            self.lecture_pieces.append({"mode": mode, "text": text.strip(), "ts": time.time()})



    def create_flashcard(self):
        if not hasattr(self, "last_extracted_text") or not self.last_extracted_text:
            QtWidgets.QMessageBox.warning(self, "No Content", "You need to summarize a slide first.")
            return

        #self.label.setText("Generating flashcards...")
        threading.Thread(target=self.generate_flashcards_from_gpt, args=(self.last_extracted_text,), daemon=True).start()



    def generate_flashcards_from_gpt(self, slide_text, extra_instruction: Optional[str] = None):

        try:
            prompt = f"""
    You are a  flashcard generator.

The following text is taken from a  lecture slide:

"\"\"{slide_text}"\"\"

Your task is to generate high-yield Q&A-style flashcards for  students. Each flashcard must target content that could realistically appear in clinical MCQs or written exams.

🎯 INSTRUCTIONS:
Extract only medically relevant, exam-appropriate material.
Focus on diagnoses, mechanisms, symptoms, treatments, first-line drugs, investigations, key cutoffs, pathways, and classic clinical signs.
Avoid trivia or background info not relevant to exams.
✅ FORMAT:
For each card, output exactly:

Question: ...
Answer: ...
Each card must:

Be concise, clear, and high-yield
Use bullet-point style in answers — short, direct phrases (not full sentences)
Follow Anki-style formatting: easy to read and memorize
❌ Do NOT:
Add explanations or teaching
Include general knowledge not found in exams
Use fluffy or vague phrasing
Include "fun facts" or contextless details
🎓 The final output must be clean, flashcard-ready, and focused entirely on what a student would need to recall under exam pressure.

Keep both question and answer very short (3–8 words).
Answers:
• Single fact → just a few words
• Multiple facts → bullet points
Focus ONLY on exam-relevant content:
• Definitions
• Classic signs & symptoms
• Cutoffs & values
• First-line investigations
• First-line treatments/drugs
• Mechanisms of action
• Key complications
NOT EVERYTHING IN A SLIDE HAS TO BE A FLASHCARD, ONLY FOCUS ON 2-4 MAIN CONCEPTS THAT CAN COME UP IN AN EXAM. DO NOT INCLUDE VERY BASIC THINGS THAT THE STUDENT SHOULD ALREADY know.
the following are a few examples: DO NOT REPRODUOCE OR OUTPUT THIS EVER IT IS JUST FOR YOU TO UNDERSTAND THE FORMATTING
Question: What is the first-line treatment for strep throat?
Answer: Penicillin V
Question what are the treatments for X?
Answer: Doxycline (first line), penicillin
Question: Where is aldosterone produced?
Answer: Zona glomerulosa (adrenal cortex)
Question: What nerve innervates the deltoid?
Answer: Axillary nerve
Question: Cutoff for diabetes diagnosis (fasting glucose)?
Answer: ≥7.0 mmol/L
Final Output:
Only clean, concise, flashcard-ready Q&A pairs. Nothing else.

    """
            
            if extra_instruction:
                prompt = (
                    "USER INSTRUCTION (highest priority): " + extra_instruction + "\n\n" + prompt
                )


            response = openai.ChatCompletion.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "You are a medical flashcard generator."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.5,
            )

            content = response['choices'][0]['message']['content']

            # Extract Q&A pairs from GPT output
            import re
            cards = re.findall(
                r'Question[:\s]+(.+?)\s+Answer[:\s]+(.+?)(?=\nQuestion[:\s]+|\Z)', 
                content, 
                re.IGNORECASE | re.DOTALL
            )

            print("🃏 RAW FLASHCARD OUTPUT 🃏")
            print(content)
            print("────────────────────────")

            count = 0
            
            




            for q, a in cards:
                self.flashcards.append((q.strip(), a.strip()))
                QtCore.QMetaObject.invokeMethod(self, "generate_list_item_safe", QtCore.Qt.ConnectionType.QueuedConnection, QtCore.Q_ARG(str, q.strip()), QtCore.Q_ARG(str, a.strip()))


     



        except Exception as e:
            print("Flashcard generation error:", e)
            #QtCore.QTimer.singleShot(0, lambda: self.label.setText("Error generating flashcards ❌"))


    @QtCore.pyqtSlot(str, str)
    def generate_list_item_safe(self, q, a):
        self.generate_list_item(q, a)



    def export_deck(self):
        if not self.flashcards:
            QtWidgets.QMessageBox.warning(self, "Nothing to Export", "No flashcards created.")
            return

        if not self.anki_filename:
            filename, _ = QtWidgets.QFileDialog.getSaveFileName(self, "Save Anki Deck", "", "TSV Files (*.tsv)")
            if not filename:
                return
            if not filename.endswith(".tsv"):
                filename += ".tsv"
            self.anki_filename = filename

        with open(self.anki_filename, 'w', newline='', encoding='utf-8') as csvfile:
            writer = csv.writer(csvfile, delimiter='\t', quotechar='"', quoting=csv.QUOTE_MINIMAL)
            for q, a in self.flashcards:
                lines = [l.strip() for l in a.splitlines() if l.strip()]
                clean_lines = [l.lstrip('–- ').strip() for l in lines]
                flat_answer = '; '.join(clean_lines)
                writer.writerow([q, flat_answer])

        QtWidgets.QMessageBox.information(self, "Exported", f"Anki deck saved: {self.anki_filename}")


    def remove_flashcard(self, question, answer):
        # Remove from internal list
        self.flashcards = [fc for fc in self.flashcards if not (fc[0] == question and fc[1] == answer)]
        
        # Refresh QListWidget
        self.flashcard_list.clear()
        for q, a in self.flashcards:
            self.generate_list_item(q, a)

    def generate_list_item(self, q, a):
        item_widget = QtWidgets.QWidget()
        layout = QtWidgets.QHBoxLayout()
        layout.setContentsMargins(5, 5, 5, 5)

        label = QtWidgets.QLabel(f"<b>Q:</b> {q}<br><b>A:</b> {a}")
        label.setWordWrap(True)
        label.setStyleSheet("font-size: 13px;")
        label.setSizePolicy(QtWidgets.QSizePolicy.Policy.Expanding, QtWidgets.QSizePolicy.Policy.Expanding)


      

        layout.addWidget(label)
      
        item_widget.setLayout(layout)

        item = QtWidgets.QListWidgetItem()
        item.setSizeHint(item_widget.sizeHint())
        self.flashcard_list.addItem(item)
        self.flashcard_list.setItemWidget(item, item_widget)

    


class QuizDialog(QtWidgets.QDialog):
    """
    Compact, scrollable quiz dialog with a Skip button.
    """
    def __init__(self, mcqs, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Akson — Quick Quiz")
        self.setMinimumWidth(300)
        self.setFixedWidth(300)
        self.setMinimumHeight(450)   # smaller than before
        self.setSizeGripEnabled(False)
        self.mcqs = mcqs
        self.i = 0
        self.wrongs = []
        self.correct = 0
        self.phase = "ask"
        self._skipped = False

        self.setStyleSheet("""
            QDialog { background:#0d0d0f; color:white; border-radius:8px; }
            QLabel  { color:white; font-size:13.5px; }
            QRadioButton { color:white; font-size:13.25px; }
            QTextBrowser { background:transparent; border:none; color:#cfcfcf; font-size:13px; }
            QPushButton {
                background: transparent; color: #dcdcdc;
                border: 0.5px solid #2a2a2e; border-radius:8px;
                padding: 6px 10px; font-weight:600; font-size:12.75px;
            }
            QPushButton:hover { background-color:  #1e1e1e; border: 0.5px solid #7f7fff; color:#d0caff; }
        """)
        self.setWindowOpacity(0.95)

        self.scroll = QtWidgets.QScrollArea()
        self.scroll.setWidgetResizable(True)
        self.scroll.setHorizontalScrollBarPolicy(QtCore.Qt.ScrollBarPolicy.ScrollBarAlwaysOff)
        self.scroll.setStyleSheet("QScrollArea { border:none; }")

        self.body = QtWidgets.QWidget()
        self.scroll.setWidget(self.body)

        self.q_label = QtWidgets.QLabel("")
        self.q_label.setWordWrap(True)
        self.q_label.setStyleSheet("font-size:14.5px; font-weight:600; color:white;")

        self.opts = [QtWidgets.QRadioButton() for _ in range(4)]
        self.group = QtWidgets.QButtonGroup(self)
        for r in self.opts:
            self.group.addButton(r)

        self.feedback = QtWidgets.QTextBrowser()
        self.feedback.hide()

        body_layout = QtWidgets.QVBoxLayout(self.body)
        body_layout.setContentsMargins(12, 12, 12, 12)
        body_layout.setSpacing(8)
        body_layout.addWidget(self.q_label)
        for r in self.opts:
            body_layout.addWidget(r)
        body_layout.addSpacing(4)
        body_layout.addWidget(self.feedback)

        self.check_btn = QtWidgets.QPushButton("Check")
        self.check_btn.clicked.connect(self._on_check_or_next)

        self.skip_btn = QtWidgets.QPushButton("Skip Quiz")
        self.skip_btn.clicked.connect(self._on_skip)

        self.close_btn = QtWidgets.QPushButton("Close")
        self.close_btn.setVisible(False)
        self.close_btn.clicked.connect(self.accept)

        footer = QtWidgets.QHBoxLayout()
        footer.addWidget(self.skip_btn)
        footer.addStretch()
        footer.addWidget(self.check_btn)
        footer.addWidget(self.close_btn)

        main = QtWidgets.QVBoxLayout(self)
        main.setContentsMargins(8, 8, 8, 8)
        main.setSpacing(6)
        main.addWidget(self.scroll)
        main.addLayout(footer)

        self._load()

    def get_wrongs(self): return self.wrongs
    def get_score(self):  return {"correct": self.correct, "total": len(self.mcqs)}
    def was_skipped(self): return self._skipped

    def _load(self):
        q = self.mcqs[self.i]
        self.q_label.setText(f"Q{self.i+1}. {q['question']}")
        for idx, r in enumerate(self.opts):
            r.setText(q['choices'][idx])
            r.setChecked(False)
            r.setStyleSheet("color:white;")
        self.feedback.hide()
        self.feedback.setHtml("")
        self.check_btn.setText("Check")
        self.close_btn.setVisible(False)
        self.phase = "ask"
        self.scroll.verticalScrollBar().setValue(0)

    def _grade(self, sel_idx):
        q = self.mcqs[self.i]
        corr = q['answer_index']
        for j, r in enumerate(self.opts):
            if j == corr:
                r.setStyleSheet("color:#9ae6b4;")
            elif j == sel_idx:
                r.setStyleSheet("color:#fca5a5;")
            else:
                r.setStyleSheet("color:#bdbdbd;")
        if sel_idx == corr:
            self.correct += 1
            self.feedback.setStyleSheet("QTextBrowser{color:#9ae6b4;}")
            self.feedback.setHtml(f"✔ Correct.<br>{q.get('explanation','')}")
        else:
            self.wrongs.append(q)
            self.feedback.setStyleSheet("QTextBrowser{color:#fca5a5;}")
            correct_txt = q['choices'][corr]
            expl = q.get("explanation","")
            self.feedback.setHtml(f"✘ Incorrect. <b>Correct:</b> {correct_txt}<br>{expl}")
        self.feedback.show()

    def _on_check_or_next(self):
        if self.phase == "ask":
            sel_idx = next((i for i, r in enumerate(self.opts) if r.isChecked()), None)
            if sel_idx is None:
                QtWidgets.QMessageBox.information(self, "Select", "Choose an option.")
                return
            self._grade(sel_idx)
            self.phase = "feedback"
            last = (self.i >= len(self.mcqs) - 1)
            self.check_btn.setText("Finish" if last else "Next")
        else:
            if self.i < len(self.mcqs) - 1:
                self.i += 1
                self._load()
            else:
                self.feedback.setStyleSheet("QTextBrowser{color:#cfcfcf;}")
                self.feedback.setHtml(f"<b>Score:</b> {self.correct}/{len(self.mcqs)}")
                self.feedback.show()
                self.close_btn.setVisible(True)
                self.check_btn.setVisible(False)

    def _on_skip(self):
        self._skipped = True
        self.reject()


# Main app loop
def window():
    # Standalone mode - no server needed
    # subprocess.Popen(["python3", "server.py"], cwd=os.path.dirname(__file__))
    app = QApplication(sys.argv)

    while True:
        # 1) show login dialog
        user_data = load_or_prompt_user()
        if user_data is None:
            sys.exit(0)

        # ---- CONSENT GATE START ----
        # Standalone mode - bypass consent
        user_data["consent_version"] = POLICY_VERSION
        # ---- CONSENT GATE END ----

        # 2) sync status and launch main window with the user_data
        sync_premium_status(user_data)
        win = MyWindow(user_data)
        win.show()

        # 3) run until close or logout
        exit_code = app.exec()
        if exit_code == 100:
            win.deleteLater()
            continue
        else:
            sys.exit(exit_code)





if __name__ == "__main__":
    window()
