"""
Akson Web Application - Full Backend API
FastAPI server with OpenAI integration for AI-powered study features.
"""

from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.responses import FileResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from pathlib import Path
from typing import Dict, Any, List, Optional
import json
import os
import re
import tempfile
import openai

# Try to import PDF/PPTX libraries
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import genanki
except ImportError:
    genanki = None

# Initialize OpenAI
try:
    from akson.config import OPENAI_API_KEY
    openai.api_key = OPENAI_API_KEY
except ImportError:
    openai.api_key = os.getenv("OPENAI_API_KEY")

app = FastAPI(title="Akson Web API", version="2.0.0")

# CORS middleware for frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

BASE_DIR = Path(__file__).resolve().parent.parent
METADATA_FILE = BASE_DIR / "metadata" / "library-metadata.json"
FRONTEND_DIR = BASE_DIR / "frontend"
PDFJS_DIR = BASE_DIR / "pdfjs"
UPLOADS_DIR = BASE_DIR / "uploads"

# Ensure uploads directory exists
UPLOADS_DIR.mkdir(exist_ok=True)


# ===================== PYDANTIC MODELS =====================

class TextRequest(BaseModel):
    text: str
    custom_prompt: Optional[str] = None


class FlashcardRequest(BaseModel):
    text: str
    extra_instruction: Optional[str] = None


class QuizRequest(BaseModel):
    text: str
    num_questions: int = 5


class ChatRequest(BaseModel):
    question: str
    context: Optional[str] = None


class SaveSessionRequest(BaseModel):
    lecture_name: str
    summaries: List[str]
    flashcards: List[Dict[str, str]]


class ExportDeckRequest(BaseModel):
    deck_name: str
    cards: List[Dict[str, str]]


# ===================== AI PROMPTS (Exact from Desktop) =====================

EXPLAIN_PROMPT = """You are Akson, a study co-pilot and educator. Your job is to convert a raw lecture slide into a high-yield, deeply explained, practical, exam-ready teaching resource. You're not summarizing — you're teaching.

🧠 CORE INSTRUCTIONS

If only one or two or a couple of words are given, then DEFINE the words, even if multiple words are given, start by shortly defining the words and where they come from. THIS IS CRUCIAL!!! After that give a line break and continue;
1️⃣ Read & Filter Smartly
Carefully read the entire slide. Completely ignore accidental or irrelevant text. Focus only on important content.

2️⃣ Teach, Don't Rephrase
Do not paraphrase or simplify blindly. Instead, explain everything clearly and intelligently — as if teaching a bright but uncertain student. For every concept, mechanism, process, term, abbreviation, or instruction:

Define it (briefly and clearly)
Explain what it does
Explain why it matters in practice
Add what the slide leaves out but a good teacher would include
Assume the student has knowledge gaps. Fill them proactively.

3️⃣ Fill In What the Slide Leaves Out
If the slide vaguely says things like "evaluation," "steps," or "solutions," you must:

List what those tests might be
Briefly explain why they're ordered
Include standard options or typical approaches where appropriate
Never allow vague or incomplete statements to pass through unexplained.
4️⃣ Always Add Real-World Context
Every explanation must tie into real-world applications. Show how the content connects to:

Key signals or indicators
Reasoning and decision steps
Investigations or tools
Solutions or approaches (baseline vs. advanced)
Outcomes
Mechanisms and potential pitfalls
5️⃣ Prioritise High-Yield Content
Mark 🔑 high-yield facts, 💡 mnemonics or tips, and ⚠️ exam traps or common misconceptions.
Mention anything that's commonly tested or clinically critical. Highlight contrasts and "commonly confused" points.

6️⃣ Use Active Recall Sparingly
After each major idea or section, insert 1 short active recall question (if appropriate).
Example: 💭 "What's the main risk factor for this condition?"
Keep questions simple and relevant.

7️⃣ Structure for Maximum Learning
Use a clean, consistent layout:

Clear section headings
Bullet points
Line spacing for readability
Concise but precise language
Avoid waffle. Use emojis (🧠, 💉, ⚠️, ❤️) sparingly to boost engagement.
8️⃣ Always Add Value Beyond the Slide
If the slide is shallow or incomplete, you must enrich it. Fill in mechanisms. Suggest simple mnemonics. Clarify unclear pathways. Break complex concepts into digestible steps. Think: "If the student sees this in a viva or ward round, what would they need to say or understand?"

9️⃣ Your Mental Model
You are a sharp, exam-aware tutor.
You are here to build true understanding, not copy text.
The final output must help the student:

Grasp the concept
Retain it
Apply it
Pass their exams with confidence

If TEXT IS SELECTED or hovered over or highlighted:
only explain the word selected, or the phrase or sentence(s) selected.


{text}
"""

SUMMARY_PROMPT = """You summarise OCR text from lecture slides into SHORT, high-yield bullets for exam revision.

GOAL
Return a compact list of independent facts in arrow/abbreviation style (memory notes), not prose.

HARD RULES (must obey)
 - start with a heading so we know what you will speak about.
 - Make it short but easy to understand. use words and wording that makes it easy to comprehend

- 4–6 bullets total. One fact per bullet. No paragraphs.
- Do NOT start bullets with category labels (e.g., "Epidemiology:", "Prevalence:", "Demographics:", "Skeletal destruction:").
  • Allowed labels only when logically required: **Criteria** and **Comparison**.
- Never chain multiple facts with ";" or long commas. Split into separate bullets.
- Use symbols & shorthand aggressively:
  ↑ increased/elevated, ↓ decreased/low, → leads to/causes, ~ approx, ≈ ratio, κ/λ light chains.
  - Prefer standard abbreviations when obvious from context.
- Ratios compact: "M:F ≈ 3:1", "Blacks:Whites ≈ 2:1".
- Minimal bold: only the key term in a bullet (e.g., **NSTEMI**, **Troponin I**, **Russell bodies**).
- Zero filler or basic definitions.

REASONING (internal; do not output)
1) Identify intent (definition/essence, criteria, comparison, atypical features, pathology/morphology, epidemiology).
2) Select the 4–6 highest-yield, exam-testable facts; prefer differentiators and criteria.
3) Rewrite each as: **Label** → compressed fact(s) with symbols/abbrevs.
4) Self-check before output:
   - Bullet count 4–6 or as little as possible.
   - No bullet starts with a category label (except **Criteria**/**Comparison** when necessary).
   - No ";" inside bullets.
   - Each bullet ≤ ~18–22 words.

FORMATS WHEN NEEDED
- **Criteria:** Use a single bullet if it's the classic triad/"≥2 of 3". Otherwise split into bullets.
- **Comparison:** Up to 3 bullets, one per entity, each starting with a bold entity then its defining line.

Return ONLY the bullets.


{text}
"""

FLASHCARD_PROMPT = """You are a flashcard generator.

The following text is taken from a lecture slide:

\"\"\"{text}\"\"\"

Your task is to generate high-yield Q&A-style flashcards for students. Each flashcard must target content that could realistically appear in clinical MCQs or written exams.

🎯 INSTRUCTIONS:
Extract only medically relevant, exam-appropriate material.
Focus on diagnoses, mechanisms, symptoms, treatments, first-line drugs, investigations, key cutoffs, pathways, and classic clinical signs.
Avoid trivia or background info not relevant to exams.
✅ FORMAT:
For each card, output exactly:

Question: ...
Answer: ...
Each card must:

Be concise, clear, and high-yield
Use bullet-point style in answers — short, direct phrases (not full sentences)
Follow Anki-style formatting: easy to read and memorize
❌ Do NOT:
Add explanations or teaching
Include general knowledge not found in exams
Use fluffy or vague phrasing
Include "fun facts" or contextless details
🎓 The final output must be clean, flashcard-ready, and focused entirely on what a student would need to recall under exam pressure.

Keep both question and answer very short (3–8 words).
Answers:
• Single fact → just a few words
• Multiple facts → bullet points
Focus ONLY on exam-relevant content:
• Definitions
• Classic signs & symptoms
• Cutoffs & values
• First-line investigations
• First-line treatments/drugs
• Mechanisms of action
• Key complications
NOT EVERYTHING IN A SLIDE HAS TO BE A FLASHCARD, ONLY FOCUS ON 2-4 MAIN CONCEPTS THAT CAN COME UP IN AN EXAM. DO NOT INCLUDE VERY BASIC THINGS THAT THE STUDENT SHOULD ALREADY know.
the following are a few examples: DO NOT REPRODUOCE OR OUTPUT THIS EVER IT IS JUST FOR YOU TO UNDERSTAND THE FORMATTING
Question: What is the first-line treatment for strep throat?
Answer: Penicillin V
Question what are the treatments for X?
Answer: Doxycline (first line), penicillin
Question: Where is aldosterone produced?
Answer: Zona glomerulosa (adrenal cortex)
Question: What nerve innervates the deltoid?
Answer: Axillary nerve
Question: Cutoff for diabetes diagnosis (fasting glucose)?
Answer: ≥7.0 mmol/L
Final Output:
Only clean, concise, flashcard-ready Q&A pairs. Nothing else.
"""

MCQ_SYSTEM_PROMPT = (
    "You write single-best-answer clinical MCQs for medical students. "
    "Use clinical vignettes and realistic distractors. Avoid trivia, avoid 'All of the above'. "
    "Make the questions like questions that would actually come up in medical school exams, completely based on the content provided. "
    "Target diagnosis, best initial investigation, best next step, or first-line management. Exactly 4 options."
)


# ===================== HELPER FUNCTIONS =====================

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF using PyMuPDF."""
    if fitz is None:
        raise HTTPException(status_code=500, detail="PyMuPDF not installed")
    
    text_parts = []
    with fitz.open(file_path) as doc:
        for page in doc:
            text_parts.append(page.get_text())
    return "\n\n".join(text_parts)


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX."""
    if Presentation is None:
        raise HTTPException(status_code=500, detail="python-pptx not installed")
    
    prs = Presentation(file_path)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_parts.append(shape.text)
    return "\n\n".join(text_parts)


def parse_flashcards(content: str) -> List[Dict[str, str]]:
    """Parse flashcard Q&A pairs from GPT output."""
    cards = re.findall(
        r'Question[:\s]+(.+?)\s+Answer[:\s]+(.+?)(?=\nQuestion[:\s]+|\Z)',
        content,
        re.IGNORECASE | re.DOTALL
    )
    return [{"question": q.strip(), "answer": a.strip()} for q, a in cards]


async def stream_openai_response(prompt: str, system_msg: str = "You are a helpful study tutor."):
    """Stream OpenAI response for real-time updates."""
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": system_msg},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            stream=True
        )
        
        for chunk in response:
            delta = chunk["choices"][0]["delta"].get("content", "")
            if delta:
                yield f"data: {json.dumps({'content': delta})}\n\n"
        
        yield "data: [DONE]\n\n"
    except Exception as e:
        yield f"data: {json.dumps({'error': str(e)})}\n\n"


# ===================== API ENDPOINTS =====================

@app.get("/health")
def health():
    """Health check endpoint."""
    return {"status": "ok", "version": "2.0.0"}


@app.get("/library")
def list_library():
    """List all saved lectures/notes."""
    if not METADATA_FILE.exists():
        return {"items": [], "summary": {}}
    with open(METADATA_FILE, "r", encoding="utf-8") as f:
        data = json.load(f)
    return {"items": data.get("files", []), "summary": data.get("summary", {})}


@app.get("/library/{note_id}")
def get_note(note_id: str):
    """Get a specific note by ID."""
    if not METADATA_FILE.exists():
        raise HTTPException(status_code=404, detail="Metadata file not found")
    with open(METADATA_FILE, "r", encoding="utf-8") as f:
        data = json.load(f)
    for note in data.get("files", []):
        if note.get("noteId") == note_id:
            return note
    raise HTTPException(status_code=404, detail="Note not found")


@app.post("/api/explain")
async def explain_text(request: TextRequest):
    """Explain selected text using the exact desktop prompt."""
    prompt = EXPLAIN_PROMPT.format(text=request.text.strip())
    
    if request.custom_prompt:
        system_msg = (
            "You are a helpful study tutor.\n\n"
            "# USER OVERRIDE INSTRUCTIONS (highest priority):\n"
            f"{request.custom_prompt}\n"
            "If any instruction above conflicts with previous instructions, "
            "FOLLOW THE USER OVERRIDE INSTRUCTIONS."
        )
    else:
        system_msg = "You are a helpful study tutor."
    
    return StreamingResponse(
        stream_openai_response(prompt, system_msg),
        media_type="text/event-stream"
    )


@app.post("/api/summarize")
async def summarize_text(request: TextRequest):
    """Summarize text using the exact desktop prompt."""
    prompt = SUMMARY_PROMPT.format(text=request.text.strip())
    
    system_msg = "You are a helpful study tutor."
    if request.custom_prompt:
        system_msg += f"\n\n# USER OVERRIDE INSTRUCTIONS:\n{request.custom_prompt}"
    
    return StreamingResponse(
        stream_openai_response(prompt, system_msg),
        media_type="text/event-stream"
    )


@app.post("/api/flashcards")
async def generate_flashcards(request: FlashcardRequest):
    """Generate flashcards from text."""
    prompt = FLASHCARD_PROMPT.format(text=request.text.strip())
    
    if request.extra_instruction:
        prompt = f"USER INSTRUCTION (highest priority): {request.extra_instruction}\n\n{prompt}"
    
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a medical flashcard generator."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.5
        )
        
        content = response['choices'][0]['message']['content']
        cards = parse_flashcards(content)
        
        return {"cards": cards, "raw": content}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/quiz")
async def generate_quiz(request: QuizRequest):
    """Generate MCQ quiz from text."""
    user_prompt = (
        f"Create {request.num_questions} clinical MCQs from the content below.\n"
        "Constraints:\n"
        "- Short vignette (1–3 sentences) when useful.\n"
        "- 4 plausible options; one correct (answer_index 0–3).\n"
        "- Use named tests/drugs and sensible cutoffs if relevant.\n"
        "- Provide a one-sentence explanation that states *why* the correct option is best.\n"
        "Output STRICT JSON (list of objects) with keys: question, choices (array of 4 strings), answer_index, explanation.\n\n"
        f"CONTENT:\n{request.text}"
    )
    
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": MCQ_SYSTEM_PROMPT},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.3
        )
        
        raw = response["choices"][0]["message"]["content"]
        
        try:
            questions = json.loads(raw)
        except json.JSONDecodeError:
            # Try to extract JSON from markdown code block
            match = re.search(r"```json\s*(.+?)\s*```", raw, re.S | re.I)
            if match:
                questions = json.loads(match.group(1))
            else:
                questions = []
        
        return {"questions": questions}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/chat")
async def chat_with_ai(request: ChatRequest):
    """Chat with AI about the current context."""
    context_part = ""
    if request.context:
        context_part = f"\n\nContext from current document:\n{request.context}\n\n"
    
    prompt = f"{context_part}Student's question: {request.question}"
    
    return StreamingResponse(
        stream_openai_response(prompt, "You are a helpful study tutor who answers questions clearly and concisely."),
        media_type="text/event-stream"
    )


@app.post("/api/upload")
async def upload_file(file: UploadFile = File(...)):
    """Upload a PDF or PPTX file and extract text."""
    # Save uploaded file
    file_path = UPLOADS_DIR / file.filename
    
    with open(file_path, "wb") as f:
        content = await file.read()
        f.write(content)
    
    # Extract text based on file type
    suffix = Path(file.filename).suffix.lower()
    
    try:
        if suffix == ".pdf":
            text = extract_text_from_pdf(str(file_path))
        elif suffix in [".pptx", ".ppt"]:
            text = extract_text_from_pptx(str(file_path))
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {suffix}")
        
        return {
            "filename": file.filename,
            "text": text,
            "path": str(file_path)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/export-anki")
async def export_anki_deck(request: ExportDeckRequest):
    """Export flashcards as Anki deck (.apkg)."""
    if genanki is None:
        raise HTTPException(status_code=500, detail="genanki not installed")
    
    import random
    
    # Create model
    model = genanki.Model(
        random.randint(1000000000, 9999999999),
        'Akson Basic Model',
        fields=[
            {'name': 'Question'},
            {'name': 'Answer'},
        ],
        templates=[
            {
                'name': 'Card 1',
                'qfmt': '{{Question}}',
                'afmt': '{{FrontSide}}<hr id="answer">{{Answer}}',
            },
        ]
    )
    
    # Create deck
    deck = genanki.Deck(
        random.randint(1000000000, 9999999999),
        request.deck_name
    )
    
    # Add cards
    for card in request.cards:
        note = genanki.Note(
            model=model,
            fields=[card['question'], card['answer']]
        )
        deck.add_note(note)
    
    # Save to temp file
    with tempfile.NamedTemporaryFile(suffix='.apkg', delete=False) as f:
        genanki.Package(deck).write_to_file(f.name)
        return FileResponse(
            f.name,
            media_type='application/octet-stream',
            filename=f"{request.deck_name}.apkg"
        )


@app.post("/api/save-session")
async def save_session(request: SaveSessionRequest):
    """Save a study session (summaries + flashcards)."""
    sessions_file = BASE_DIR / "sessions.json"
    
    if sessions_file.exists():
        with open(sessions_file, "r") as f:
            sessions = json.load(f)
    else:
        sessions = {}
    
    sessions[request.lecture_name] = {
        "summaries": request.summaries,
        "flashcards": request.flashcards,
        "created": str(Path().stat().st_mtime) if Path().exists() else "now"
    }
    
    with open(sessions_file, "w") as f:
        json.dump(sessions, f, indent=2)
    
    return {"status": "saved", "lecture": request.lecture_name}


@app.get("/api/sessions")
def list_sessions():
    """List all saved sessions."""
    sessions_file = BASE_DIR / "sessions.json"
    
    if not sessions_file.exists():
        return {"sessions": {}}
    
    with open(sessions_file, "r") as f:
        return {"sessions": json.load(f)}


# ===================== STATIC FILES =====================

@app.get("/")
def root():
    """Serve the frontend."""
    index = FRONTEND_DIR / "index.html"
    if index.exists():
        return FileResponse(index)
    raise HTTPException(status_code=404, detail="Frontend not built")


# Mount static directories
if FRONTEND_DIR.exists():
    app.mount("/static", StaticFiles(directory=FRONTEND_DIR), name="static")

if PDFJS_DIR.exists():
    app.mount("/pdfjs", StaticFiles(directory=PDFJS_DIR), name="pdfjs")

# Serve uploaded files
if UPLOADS_DIR.exists():
    app.mount("/uploads", StaticFiles(directory=UPLOADS_DIR), name="uploads")
